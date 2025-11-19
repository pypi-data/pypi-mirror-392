"""PPTX/PPT parsing utilities using officeparser."""

from __future__ import annotations


async def parse_pptx(file_path: str) -> dict:
    if not file_path:
        return {"text": ""}
    
    # Python doesn't have officeparser, use python-pptx instead
    from pptx import Presentation
    
    prs = Presentation(file_path)
    texts = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    
    return {"text": "\n\n".join(texts)}
