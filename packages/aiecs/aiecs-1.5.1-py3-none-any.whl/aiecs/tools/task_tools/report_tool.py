"""
Report Tool - A multi-format report generation tool supporting HTML, PDF, Excel, PowerPoint, Markdown, Word, and image-based reports.

This module provides a comprehensive report generation tool that can be used to create various types of reports
in different formats. It supports template-based rendering, data visualization, and batch processing.

Author: Your Organization
Version: 1.0.0
"""

import os
import bleach
from typing import Dict, Any, List, Optional, Union, Tuple, Set
from jinja2 import FileSystemLoader, sandbox

# from weasyprint import HTML  # TODO: Re-enable when deployment issues
# are resolved
import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from docx import Document
from docx.shared import Pt as DocxPt
import matplotlib.pyplot as plt
from pydantic import BaseModel, ConfigDict, Field
import tempfile
import logging

from aiecs.tools.base_tool import BaseTool
from aiecs.tools import register_tool
from aiecs.tools.temp_file_manager import TempFileManager


# Exceptions
class ReportToolError(Exception):
    """Base exception for ReportTool errors."""


class FileOperationError(ReportToolError):
    """Raised when file operations fail."""


# Helper function for HTML sanitization


def sanitize_html(
    html_content: str,
    allowed_tags: Set[str],
    allowed_attributes: Dict[str, List[str]],
) -> str:
    """
    Sanitize HTML content to prevent XSS attacks.

    Args:
        html_content (str): The HTML content to sanitize.
        allowed_tags (Set[str]): Set of allowed HTML tags.
        allowed_attributes (Dict[str, List[str]]): Dictionary of allowed attributes for each tag.

    Returns:
        str: Sanitized HTML content.
    """
    return bleach.clean(
        html_content,
        tags=allowed_tags,
        attributes=allowed_attributes,
        strip=True,
    )


# Type alias for dataset entries
DatasetType = Dict[str, Union[pd.DataFrame, List[Dict[str, Any]]]]


@register_tool("report")
class ReportTool(BaseTool):
    """
    Multi-format report generation tool supporting HTML, Excel, PowerPoint, Markdown, Word, and image-based reports.

    NOTE: PDF generation is temporarily disabled due to weasyprint deployment complexity.

    Operations:
      - generate_html: Render HTML report using Jinja2.
      - generate_pdf: Currently disabled - will be re-enabled in future release.
      - generate_excel: Create Excel workbook with multiple sheets and styling.
      - generate_pptx: Create PowerPoint presentation with customizable slides.
      - generate_markdown: Render Markdown report using Jinja2.
      - generate_word: Create Word document with customizable styles.
      - generate_image: Generate charts (bar, line, pie) using Matplotlib.
      - batch_generate: Generate multiple reports in parallel.

    Inherits from BaseTool.
    """

    # Configuration schema
    class Config(BaseModel):
        """Configuration for the report tool"""

        model_config = ConfigDict(env_prefix="REPORT_TOOL_")

        templates_dir: str = Field(
            default=os.getcwd(), description="Directory for Jinja2 templates"
        )
        default_output_dir: str = Field(
            default=os.path.join(tempfile.gettempdir(), "reports"),
            description="Default directory for output files",
        )
        allowed_extensions: List[str] = Field(
            default=[
                ".html",
                ".pdf",
                ".xlsx",
                ".pptx",
                ".docx",
                ".md",
                ".png",
            ],
            description="Allowed file extensions for outputs",
        )
        pdf_page_size: str = Field(default="A4", description="Default PDF page size")
        default_font: str = Field(default="Arial", description="Default font for documents")
        default_font_size: int = Field(default=12, description="Default font size in points")
        allowed_html_tags: Set[str] = Field(
            default={
                "h1",
                "h2",
                "h3",
                "h4",
                "h5",
                "h6",
                "p",
                "br",
                "a",
                "ul",
                "ol",
                "li",
                "strong",
                "em",
                "b",
                "i",
                "table",
                "tr",
                "td",
                "th",
                "thead",
                "tbody",
                "span",
                "div",
                "img",
                "hr",
                "code",
                "pre",
            },
            description="Allowed HTML tags for sanitization",
        )
        allowed_html_attributes: Dict[str, List[str]] = Field(
            default={
                "a": ["href", "title", "target"],
                "img": ["src", "alt", "title", "width", "height"],
                "td": ["colspan", "rowspan", "align"],
                "th": ["colspan", "rowspan", "align"],
                "*": ["class", "id", "style"],
            },
            description="Allowed HTML attributes for sanitization",
        )
        temp_files_max_age: int = Field(
            default=3600,
            description="Maximum age of temporary files in seconds",
        )

    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        Initialize ReportTool with settings and resources.

        Args:
            config (Dict, optional): Configuration overrides for ReportTool.

        Raises:
            ValueError: If config contains invalid settings.
        """
        super().__init__(config)

        # Parse configuration
        self.config = self.Config(**(config or {}))

        self.logger = logging.getLogger(__name__)
        if not self.logger.handlers:
            handler = logging.StreamHandler()
            handler.setFormatter(logging.Formatter("%(asctime)s %(levelname)s %(message)s"))
            self.logger.addHandler(handler)
        self.logger.setLevel(logging.INFO)
        self._jinja_env = sandbox.SandboxedEnvironment(
            loader=FileSystemLoader(self.config.templates_dir), autoescape=True
        )
        self._temp_manager = TempFileManager(
            self.config.default_output_dir, self.config.temp_files_max_age
        )

    def generate_html(
        self,
        template_path: Optional[str],
        template_str: Optional[str],
        context: Dict[str, Any],
        output_path: str,
        template_variables: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Render an HTML report using a Jinja2 template.

        Args:
            template_path (Optional[str]): Path to the template file.
            template_str (Optional[str]): Template string content.
            context (Dict[str, Any]): Template context data.
            output_path (str): Path to save the HTML file.
            template_variables (Optional[Dict[str, str]]): Variables for dynamic output path.

        Returns:
            str: Path to the generated HTML file.

        Raises:
            FileOperationError: If template file is not found or writing fails.
        """
        try:
            if template_path:
                os.path.join(self.config.templates_dir, template_path)
                tmpl = self._jinja_env.get_template(template_path)
            else:
                tmpl = self._jinja_env.from_string(template_str)
            html = tmpl.render(**context)
            csrf_meta = "<meta http-equiv=\"Content-Security-Policy\" content=\"default-src 'self'; script-src 'self'; object-src 'none'\">\n"
            csrf_meta += '<meta name="referrer" content="no-referrer">\n'
            if "<head>" in html:
                html = html.replace("<head>", "<head>\n" + csrf_meta)
            else:
                html = csrf_meta + html
            html = sanitize_html(
                html,
                self.config.allowed_html_tags,
                self.config.allowed_html_attributes,
            )
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(html)
            self._temp_manager.register_file(output_path)
            return output_path
        except Exception as e:
            raise FileOperationError(f"Failed to generate HTML: {str(e)}")

    def generate_pdf(
        self,
        html: Optional[str],
        html_schema: Optional[Dict],
        output_path: str,
        page_size: Optional[str] = None,
        template_variables: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Generate a PDF report from HTML content or a Jinja2 template.

        NOTE: PDF generation is currently disabled due to weasyprint deployment complexity.
        This feature will be re-enabled in a future release.

        Args:
            html (Optional[str]): HTML content.
            html_schema (Optional[Dict]): Dict for HTML generation.
            output_path (str): Path to save the PDF file.
            page_size (Optional[str]): PDF page size.
            template_variables (Optional[Dict[str, str]]): Variables for dynamic output path.

        Returns:
            str: Path to the generated PDF file.

        Raises:
            FileOperationError: PDF generation is currently disabled.
        """
        raise FileOperationError(
            "PDF generation is currently disabled due to weasyprint deployment complexity. "
            "Please use generate_html() to create HTML reports instead. "
            "PDF functionality will be restored in a future release."
        )

        # TODO: Re-enable when weasyprint deployment issues are resolved
        # try:
        #     if not html and html_schema:
        #         html_path = self.generate_html(**html_schema)
        #         with open(html_path, 'r', encoding='utf-8') as f:
        #             html = f.read()
        #     HTML(string=html).write_pdf(
        #         output_path,
        #         stylesheets=[{'page_size': page_size or self.settings.pdf_page_size}]
        #     )
        #     self._temp_manager.register_file(output_path)
        #     return output_path
        # except Exception as e:
        #     raise FileOperationError(f"Failed to generate PDF: {str(e)}")

    def generate_excel(
        self,
        sheets: Dict[str, Union[pd.DataFrame, List[Dict[str, Any]]]],
        output_path: str,
        styles: Optional[Dict[str, Dict[str, Any]]] = None,
        template_variables: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Generate an Excel workbook with multiple sheets and optional styling.

        Args:
            sheets (Dict[str, Union[pd.DataFrame, List[Dict[str, Any]]]]): Sheet data.
            output_path (str): Path to save the Excel file.
            styles (Optional[Dict[str, Dict[str, Any]]]): Cell styling.
            template_variables (Optional[Dict[str, str]]): Variables for dynamic output path.

        Returns:
            str: Path to the generated Excel file.

        Raises:
            FileOperationError: If Excel generation fails.
        """
        try:
            writer = pd.ExcelWriter(output_path, engine="xlsxwriter")
            workbook = writer.book
            for name, data in sheets.items():
                df = data if isinstance(data, pd.DataFrame) else pd.DataFrame(data)
                df.to_excel(writer, sheet_name=name[:31], index=False)
                if styles and name in styles:
                    worksheet = writer.sheets[name[:31]]
                    for cell, style in styles[name].items():
                        format_dict = {}
                        if style.get("bold"):
                            format_dict["bold"] = True
                        if style.get("font_size"):
                            format_dict["font_size"] = style["font_size"]
                        if style.get("bg_color"):
                            format_dict["bg_color"] = style["bg_color"]
                        worksheet.write(
                            cell,
                            df.loc[int(cell[1:]) - 1, cell[0]],
                            workbook.add_format(format_dict),
                        )
            writer.close()
            self._temp_manager.register_file(output_path)
            return output_path
        except Exception as e:
            raise FileOperationError(f"Failed to generate Excel: {str(e)}")

    def generate_pptx(
        self,
        slides: List[Dict],
        output_path: str,
        default_font: Optional[str] = None,
        default_font_size: Optional[int] = None,
        default_font_color: Optional[Tuple[int, int, int]] = None,
        template_variables: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Generate a PowerPoint presentation with customizable slides.

        Args:
            slides (List[Dict]): List of slide data.
            output_path (str): Path to save the PPTX file.
            default_font (Optional[str]): Default font for slides.
            default_font_size (Optional[int]): Default font size.
            default_font_color (Optional[Tuple[int, int, int]]): Default font color.
            template_variables (Optional[Dict[str, str]]): Variables for dynamic output path.

        Returns:
            str: Path to the generated PPTX file.

        Raises:
            FileOperationError: If PPTX generation fails.
        """
        try:
            prs = Presentation()
            for slide in slides:
                s = prs.slides.add_slide(prs.slide_layouts[1])
                title_shape = s.shapes.title
                title_shape.text = slide["title"]
                font = slide.get("font") or default_font or self.config.default_font
                font_size = (
                    slide.get("font_size") or default_font_size or self.config.default_font_size
                )
                slide.get("font_color") or default_font_color or (0, 0, 0)
                title_shape.text_frame.paragraphs[0].font.name = font
                title_shape.text_frame.paragraphs[0].font.size = Pt(font_size)
                # Set font color safely - skip color setting for now to avoid library issues
                # Font color setting in python-pptx can be problematic,
                # focusing on core functionality
                body = s.shapes.placeholders[1].text_frame
                for bullet in slide["bullets"]:
                    p = body.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.name = font
                    p.font.size = Pt(font_size)
                    # Skip font color setting for bullet points to avoid
                    # library issues
            prs.save(output_path)
            self._temp_manager.register_file(output_path)
            return output_path
        except Exception as e:
            raise FileOperationError(f"Failed to generate PPTX: {str(e)}")

    def generate_markdown(
        self,
        template_path: Optional[str],
        template_str: Optional[str],
        context: Dict[str, Any],
        output_path: str,
        template_variables: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Render a Markdown report using a Jinja2 template.

        Args:
            template_path (Optional[str]): Path to the template file.
            template_str (Optional[str]): Template string content.
            context (Dict[str, Any]): Template context data.
            output_path (str): Path to save the Markdown file.
            template_variables (Optional[Dict[str, str]]): Variables for dynamic output path.

        Returns:
            str: Path to the generated Markdown file.

        Raises:
            FileOperationError: If rendering or writing fails.
        """
        try:
            if template_path:
                tmpl = self._jinja_env.get_template(template_path)
            else:
                tmpl = self._jinja_env.from_string(template_str)
            markdown_content = tmpl.render(**context)
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(markdown_content)
            self._temp_manager.register_file(output_path)
            return output_path
        except Exception as e:
            raise FileOperationError(f"Failed to generate Markdown: {str(e)}")

    def generate_word(
        self,
        template_path: Optional[str],
        template_str: Optional[str],
        context: Dict[str, Any],
        output_path: str,
        font: Optional[str] = None,
        font_size: Optional[int] = None,
        font_color: Optional[Tuple[int, int, int]] = None,
        template_variables: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Generate a Word document from a Jinja2 template with customizable styles.

        Args:
            template_path (Optional[str]): Path to the template file.
            template_str (Optional[str]): Template string content.
            context (Dict[str, Any]): Template context data.
            output_path (str): Path to save the DOCX file.
            font (Optional[str]): Font for the document.
            font_size (Optional[int]): Font size.
            font_color (Optional[Tuple[int, int, int]]): Font color.
            template_variables (Optional[Dict[str, str]]): Variables for dynamic output path.

        Returns:
            str: Path to the generated DOCX file.

        Raises:
            FileOperationError: If Word generation fails.
        """
        try:
            if template_path:
                tmpl = self._jinja_env.get_template(template_path)
            else:
                tmpl = self._jinja_env.from_string(template_str)
            content = tmpl.render(**context)
            doc = Document()
            font = font or self.config.default_font
            font_size = font_size or self.config.default_font_size
            font_color = font_color or (0, 0, 0)
            for line in content.splitlines():
                p = doc.add_paragraph()
                run = p.add_run(line)
                run.font.name = font
                run.font.size = DocxPt(font_size)
                # Skip font color setting for Word documents to avoid library
                # issues
            doc.save(output_path)
            self._temp_manager.register_file(output_path)
            return output_path
        except Exception as e:
            raise FileOperationError(f"Failed to generate Word: {str(e)}")

    def generate_image(
        self,
        chart_type: str,
        data: Union[pd.DataFrame, List[Dict[str, Any]]],
        output_path: str,
        x_col: Optional[str] = None,
        y_col: Optional[str] = None,
        labels: Optional[List[str]] = None,
        title: Optional[str] = None,
        width: int = 8,
        height: int = 6,
        template_variables: Optional[Dict[str, str]] = None,
    ) -> str:
        """
        Generate a chart (bar, line, pie) using Matplotlib.

        Args:
            chart_type (str): Type of chart ('bar', 'line', 'pie').
            data (Union[pd.DataFrame, List[Dict[str, Any]]]): Chart data.
            output_path (str): Path to save the image file.
            x_col (Optional[str]): X-axis column name.
            y_col (Optional[str]): Y-axis column name.
            labels (Optional[List[str]]): Labels for pie chart.
            title (Optional[str]): Chart title.
            width (int): Chart width.
            height (int): Chart height.
            template_variables (Optional[Dict[str, str]]): Variables for dynamic output path.

        Returns:
            str: Path to the generated image file.

        Raises:
            FileOperationError: If chart generation fails.
        """
        try:
            df = data if isinstance(data, pd.DataFrame) else pd.DataFrame(data)
            plt.figure(figsize=(width, height))
            if chart_type == "bar":
                df.plot.bar(x=x_col, y=y_col, title=title)
            elif chart_type == "line":
                df.plot.line(x=x_col, y=y_col, title=title)
            elif chart_type == "pie":
                plt.pie(
                    df[y_col],
                    labels=df[x_col] if x_col else labels,
                    autopct="%1.1f%%",
                )
                plt.title(title)
            plt.savefig(output_path)
            plt.close()
            self._temp_manager.register_file(output_path)
            return output_path
        except Exception as e:
            raise FileOperationError(f"Failed to generate image: {str(e)}")

    def batch_generate(
        self,
        operation: str,
        contexts: List[Dict[str, Any]],
        output_paths: List[str],
        datasets: Optional[List[DatasetType]] = None,
        slides: Optional[List[List[Dict]]] = None,
    ) -> List[str]:
        """
        Generate multiple reports in parallel for different contexts or datasets.

        Args:
            operation (str): Operation to perform.
            contexts (List[Dict[str, Any]]): Contexts for HTML, Markdown, Word, PDF.
            output_paths (List[str]): Paths for generated files.
            datasets (Optional[List[DatasetType]]): Datasets for Excel, Image.
            slides (Optional[List[List[Dict]]]): Slides for PPTX.

        Returns:
            List[str]: List of generated file paths.

        Raises:
            FileOperationError: If batch generation fails.
        """
        try:
            tasks = []
            input_data = contexts or datasets or slides
            for i, output_path in enumerate(output_paths):
                op_params = {"output_path": output_path}
                if operation in (
                    "generate_html",
                    "generate_markdown",
                    "generate_word",
                ):
                    op_params.update(input_data[i])
                    op_params["template_path"] = input_data[i].get("template_path")
                    op_params["template_str"] = input_data[i].get("template_str")
                    if operation == "generate_word":
                        op_params["font"] = input_data[i].get("font")
                        op_params["font_size"] = input_data[i].get("font_size")
                        op_params["font_color"] = input_data[i].get("font_color")
                elif operation == "generate_excel":
                    op_params["sheets"] = input_data[i]
                    op_params["styles"] = input_data[i].get("styles")
                elif operation == "generate_pptx":
                    op_params["slides"] = input_data[i]
                    op_params["default_font"] = (
                        input_data[i][0].get("font") if input_data[i] else None
                    )
                    op_params["default_font_size"] = (
                        input_data[i][0].get("font_size") if input_data[i] else None
                    )
                    op_params["default_font_color"] = (
                        input_data[i][0].get("font_color") if input_data[i] else None
                    )
                elif operation == "generate_image":
                    op_params.update(input_data[i])
                elif operation == "generate_pdf":
                    op_params["html"] = input_data[i].get("html")
                    op_params["html_schema"] = (
                        input_data[i] if input_data[i].get("context") else None
                    )
                    op_params["page_size"] = input_data[i].get("page_size")
                tasks.append({"op": operation, "kwargs": op_params})
            # Execute tasks synchronously for batch generation
            results = []
            for task in tasks:
                op_name = task["op"]
                kwargs = task["kwargs"]

                if op_name == "generate_html":
                    result = self.generate_html(**kwargs)
                elif op_name == "generate_excel":
                    result = self.generate_excel(**kwargs)
                elif op_name == "generate_pptx":
                    result = self.generate_pptx(**kwargs)
                elif op_name == "generate_markdown":
                    result = self.generate_markdown(**kwargs)
                elif op_name == "generate_word":
                    result = self.generate_word(**kwargs)
                elif op_name == "generate_image":
                    result = self.generate_image(**kwargs)
                elif op_name == "generate_pdf":
                    result = self.generate_pdf(**kwargs)
                else:
                    raise FileOperationError(f"Unsupported operation: {op_name}")

                results.append(result)
            return results
        except Exception as e:
            raise FileOperationError(f"Failed to generate batch reports: {str(e)}")
