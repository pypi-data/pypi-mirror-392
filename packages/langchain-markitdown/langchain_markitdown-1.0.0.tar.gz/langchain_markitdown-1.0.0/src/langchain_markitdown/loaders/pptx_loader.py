from typing import TYPE_CHECKING, Dict, Any, List, Optional
from langchain_core.documents import Document
from ..base_loader import BaseMarkitdownLoader
from langchain_core.language_models import BaseChatModel
import re
import os
import io
import logging
from ..utils import langchain_caption_adapter, get_image_format  # Import both at top

if TYPE_CHECKING:
    from markitdown import MarkItDown


class PptxLoader(BaseMarkitdownLoader):
    def __init__(
        self,
        file_path: str,
        split_by_page: bool = False,
        llm: Optional[BaseChatModel] = None,
        prompt: Optional[str] = None,
        verbose: Optional[bool] = None,
        *,
        converter: Optional["MarkItDown"] = None,
    ):
        super().__init__(file_path, verbose=bool(verbose), converter=converter)
        self.split_by_page = split_by_page
        self.llm = llm
        self.prompt = prompt
        self.logger.info(f"Initialized PptxLoader for {file_path} with split_by_page={split_by_page}")
        self.logger.info(f"Langchain LLM for image captioning: {llm.__class__.__name__ if llm else 'None'}")

    def _extract_metadata(self) -> Dict[str, Any]:
        from pptx import Presentation

        metadata = {
            "source": self.file_path,
            "file_name": self._get_file_name(self.file_path),
            "file_size": self._get_file_size(self.file_path),
            "success": False,
            "conversion_success": False,
        }

        try:
            prs = Presentation(self.file_path)
            metadata["slide_count"] = len(prs.slides)
            self.logger.info(f"Found {metadata['slide_count']} slides in the presentation")

            core_props = prs.core_properties

            def add_if_present(attr_name, key=None):
                value = getattr(core_props, attr_name, None)
                if value:
                    metadata[key or attr_name] = str(value)

            for attr in ["author", "title", "subject", "keywords", "created", "modified", "last_modified_by", "category", "revision"]:
                add_if_present(attr)

            image_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 13)
            text_box_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 17)
            chart_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 3)
            table_count = sum(1 for slide in prs.slides for shape in slide.shapes if shape.shape_type == 19)

            metadata.update({
                "image_count": image_count,
                "text_box_count": text_box_count,
                "chart_count": chart_count,
                "table_count": table_count
            })

        except Exception as e:
            self.logger.warning(f"Failed to extract detailed metadata: {str(e)}")
            metadata["metadata_extraction_error"] = str(e)

        return metadata

    def _caption_images(self, markdown_content: str) -> str:
        from pptx import Presentation

        prs = Presentation(self.file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type') and shape.shape_type == 13:
                    image_data = shape.image.blob
                    stream_info = type("DummyStreamInfo", (object,), {
                        "mimetype": "", "extension": "", "name": shape.name
                    })()

                    stream_info.mimetype, stream_info.extension = get_image_format(image_data)
                    if stream_info.mimetype not in ["image/png", "image/jpeg", "image/gif", "image/webp"]:
                        self.logger.warning(f"Skipping captioning for unsupported image format: {stream_info.mimetype}")
                        continue

                    image_stream = io.BytesIO(image_data)
                    try:
                        caption = langchain_caption_adapter(
                            file_stream=image_stream,
                            stream_info=stream_info,
                            client=self.llm,
                            model=None,
                            prompt=self.prompt
                        )
                        if caption:
                            self.logger.info(f"Generated caption: {caption[:50]}...")
                            cleaned_shape_name = re.sub(r"\W", "", shape.name)
                            alt_text_pattern = r"!\[.*?]\(" + re.escape(f"{cleaned_shape_name}.jpg") + r"\)"
                            markdown_content = re.sub(alt_text_pattern, f"![{caption}]({cleaned_shape_name}.jpg)", markdown_content)
                    except Exception as e:
                        self.logger.error(f"Error during LLM captioning: {e}")
        return markdown_content

    def _split_markdown_into_documents(self, markdown_content: str, metadata: Dict[str, Any]) -> List[Document]:
        slide_pattern = r"^\n*<!-- Slide number: (\d+) -->\n"
        slide_splits = re.split(slide_pattern, markdown_content, flags=re.MULTILINE)

        documents = []
        current_page_content = slide_splits[0]
        current_page_num = 1

        for i in range(1, len(slide_splits), 2):
            if current_page_content.strip():
                page_metadata = metadata.copy()
                page_metadata.update({"page_number": current_page_num, "content_type": "presentation_slide"})
                documents.append(Document(page_content=current_page_content, metadata=page_metadata))

            current_page_num = int(slide_splits[i])
            current_page_content = slide_splits[i + 1]

        if current_page_content.strip():
            page_metadata = metadata.copy()
            page_metadata.update({"page_number": current_page_num, "content_type": "presentation_slide"})
            documents.append(Document(page_content=current_page_content, metadata=page_metadata))

        return documents if documents else [Document(page_content="", metadata=metadata)]

    def load(self, headers_to_split_on: Optional[List[str]] = None) -> List[Document]:
        self.logger.info(f"Starting to load PPTX file: {self.file_path}")
        metadata = self._extract_metadata()

        try:
            self.logger.info("Converting PPTX to markdown")
            result = self._convert_to_markdown()
        except FileNotFoundError as exc:
            metadata["error"] = "File not found."
            raise ValueError(
                f"Markitdown conversion failed for {self.file_path}: File not found",
            ) from exc
        except Exception as exc:
            metadata["error"] = str(exc)
            raise ValueError(
                f"Failed to load and convert PPTX file: {exc}",
            ) from exc

        metadata["success"] = True
        metadata["conversion_success"] = True
        metadata.update(self._extract_conversion_metadata(result))
        markdown_content = self._get_text_content(result)

        if self.llm:
            self.logger.info("Processing images and generating captions...")
            markdown_content = self._caption_images(markdown_content)

        self.logger.info(
            "Conversion complete, markdown content length: %s characters",
            len(markdown_content),
        )

        if not self.split_by_page:
            metadata["content_type"] = "presentation_full"
            return [Document(page_content=markdown_content, metadata=metadata)]

        return self._split_markdown_into_documents(markdown_content, metadata)
