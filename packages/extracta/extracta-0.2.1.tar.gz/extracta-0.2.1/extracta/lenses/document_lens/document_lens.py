import io
import json
from pathlib import Path
from typing import Dict, Any

from ..base_lens import BaseLens

# Import libraries for different file types
try:
    import PyPDF2
    from PyPDF2 import PdfReader
except ImportError:
    PyPDF2 = None  # type: ignore[assignment]
    PdfReader = None  # type: ignore[misc,assignment]

try:
    from pdfplumber.pdf import PDF as PDFPlumberPDF  # noqa: N811
except ImportError:
    PDFPlumberPDF = None  # type: ignore[assignment, misc]

try:
    from docx import Document
except ImportError:
    Document = None

# Office document support - enhanced from mark-mate
try:
    from pptx import Presentation

    pptx_available = True
except ImportError:
    Presentation = None
    pptx_available = False

try:
    import pandas as pd

    pandas_available = True
except ImportError:
    pd = None
    pandas_available = False

try:
    from openpyxl import load_workbook

    openpyxl_available = True
except ImportError:
    load_workbook = None
    openpyxl_available = False

# PPTX support removed - moved to presentation-lens
# See: https://github.com/michael-borck/presentation-lens


class DocumentLens(BaseLens):
    """Lens for extracting content from document files"""

    SUPPORTED_EXTENSIONS = {
        ".txt",
        ".md",
        ".rst",
        ".json",  # Text formats
        ".pdf",  # PDF documents
        ".docx",  # Word documents
        ".pptx",  # PowerPoint presentations
        ".xlsx",  # Excel spreadsheets
        ".xls",  # Excel spreadsheets (older format)
        ".csv",  # CSV files
        ".tsv",  # TSV files
        # PPTX moved to presentation-lens for specialized handling
    }

    def __init__(self):
        self.pdf_available = PdfReader is not None
        self.docx_available = Document is not None
        self.pptx_available = Presentation is not None
        self.excel_available = pandas_available or openpyxl_available

    def extract(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from document file"""
        try:
            # Ensure file_path is a Path object
            if isinstance(file_path, str):
                file_path = Path(file_path)

            if file_path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
                return {
                    "success": False,
                    "error": f"Unsupported file type: {file_path.suffix}",
                    "data": {},
                }

            # Route to appropriate extraction method
            if file_path.suffix.lower() == ".pdf":
                return self._extract_pdf(file_path)
            elif file_path.suffix.lower() == ".docx":
                return self._extract_docx(file_path)
            elif file_path.suffix.lower() in [".pptx"]:
                return self._extract_powerpoint(file_path)
            elif file_path.suffix.lower() in [".xlsx", ".xls"]:
                return self._extract_excel(file_path)
            elif file_path.suffix.lower() in [".csv", ".tsv"]:
                return self._extract_csv(file_path)
            elif file_path.suffix.lower() == ".json":
                return self._extract_json(file_path)
            else:
                # Text files (.txt, .md, .rst)
                return self._extract_text(file_path)

        except Exception as e:
            return {"success": False, "error": str(e), "data": {}}

    def _extract_text(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from plain text/markdown files"""
        try:
            # Try different encodings
            encodings = ["utf-8", "utf-16", "latin-1", "cp1252"]

            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as f:
                        content = f.read()
                    return {
                        "success": True,
                        "data": {
                            "content_type": "text",
                            "raw_content": content,
                            "file_path": str(file_path),
                            "file_size": file_path.stat().st_size,
                            "encoding": encoding,
                        },
                    }
                except UnicodeDecodeError:
                    continue

            # If all encodings fail, use utf-8 with error handling
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                content = f.read()
            return {
                "success": True,
                "data": {
                    "content_type": "text",
                    "raw_content": content,
                    "file_path": str(file_path),
                    "file_size": file_path.stat().st_size,
                    "encoding": "utf-8 (with error replacement)",
                },
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to read text file: {e}",
                "data": {},
            }

    def _extract_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from PDF file"""
        if not self.pdf_available:
            return {
                "success": False,
                "error": "PDF extraction not available. Install with: pip install PyPDF2",
                "data": {},
            }

        try:
            text_parts = []

            # Try with PyPDF2 first
            with open(file_path, "rb") as f:
                pdf_reader = PyPDF2.PdfReader(f)

                for page in pdf_reader.pages:
                    text = page.extract_text()
                    if text.strip():
                        text_parts.append(text)

            # If PyPDF2 didn't extract much text, try pdfplumber
            if len("".join(text_parts)) < 100 and PDFPlumberPDF is not None:
                with PDFPlumberPDF.open(file_path) as pdf:
                    text_parts = []
                    for page in pdf.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)

            if not text_parts:
                return {
                    "success": False,
                    "error": "No text could be extracted from PDF",
                    "data": {},
                }

            content = "\n\n".join(text_parts)

            data = {
                "content_type": "text",
                "raw_content": content,
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size,
                "pages": len(pdf_reader.pages) if "pdf_reader" in locals() else None,
                "extraction_method": "pdf",
            }

            # Apply cascading analysis for code detection
            from .. import analyze_extracted_content

            enhanced_data = analyze_extracted_content(data)
            if "cascading_analysis" in enhanced_data:
                data["cascading_analysis"] = enhanced_data["cascading_analysis"]

            return {
                "success": True,
                "data": data,
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract PDF: {e}",
                "data": {},
            }

    def _extract_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from DOCX file"""
        if not self.docx_available:
            return {
                "success": False,
                "error": "DOCX extraction not available. Install with: pip install python-docx",
                "data": {},
            }

        try:
            doc = Document(file_path)
            text_parts = []

            # Extract text from paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text)

            if not text_parts:
                return {
                    "success": False,
                    "error": "No text could be extracted from DOCX",
                    "data": {},
                }

            content = "\n\n".join(text_parts)
            return {
                "success": True,
                "data": {
                    "content_type": "text",
                    "raw_content": content,
                    "file_path": str(file_path),
                    "file_size": file_path.stat().st_size,
                    "paragraphs": len(doc.paragraphs),
                    "tables": len(doc.tables),
                    "extraction_method": "docx",
                },
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract DOCX: {e}",
                "data": {},
            }

    def _extract_powerpoint(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PowerPoint presentation (enhanced from mark-mate)"""
        if not self.pptx_available:
            return {
                "success": False,
                "error": "PowerPoint extraction not available. Install with: pip install python-pptx",
                "data": {},
            }

        try:
            presentation = Presentation(file_path)
            slides_content = []
            total_text = ""

            for i, slide in enumerate(presentation.slides, 1):
                slide_data = {
                    "slide_number": i,
                    "title": "",
                    "content": "",
                    "notes": "",
                }

                # Extract title and content from shapes
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        slide_text_parts.append(text)

                        # Try to identify title (usually first text box or larger font)
                        if not slide_data["title"] and len(text) < 100:
                            slide_data["title"] = text
                        else:
                            slide_data["content"] += text + "\n"

                # Extract speaker notes
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_data["notes"] = (
                        slide.notes_slide.notes_text_frame.text.strip()
                    )

                # Combine all text for this slide
                all_slide_text = "\n".join(slide_text_parts)
                if slide_data["notes"]:
                    all_slide_text += "\n[SPEAKER NOTES]\n" + slide_data["notes"]

                slides_content.append(slide_data)
                total_text += f"\n=== Slide {i} ===\n" + all_slide_text + "\n"

            # Create analysis
            analysis = {
                "total_slides": len(presentation.slides),
                "slides_with_notes": sum(1 for s in slides_content if s["notes"]),
                "avg_content_length": sum(len(s["content"]) for s in slides_content)
                / len(slides_content)
                if slides_content
                else 0,
                "slides_breakdown": slides_content[
                    :10
                ],  # First 10 slides for detailed analysis
            }

            # Create content summary
            content = f"""POWERPOINT PRESENTATION SUMMARY:
Total Slides: {analysis["total_slides"]}
Slides with Speaker Notes: {analysis["slides_with_notes"]}
Average Content Length: {analysis["avg_content_length"]:.1f} characters

PRESENTATION CONTENT:
{total_text}"""

            return {
                "success": True,
                "data": {
                    "content_type": "presentation",
                    "raw_content": content,
                    "file_path": str(file_path),
                    "file_size": file_path.stat().st_size,
                    "extraction_method": "pptx",
                    "presentation_analysis": analysis,
                },
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract PowerPoint: {e}",
                "data": {},
            }

    def _extract_excel(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from Excel workbook (enhanced from mark-mate)"""
        if not self.excel_available:
            return {
                "success": False,
                "error": "Excel extraction not available. Install with: pip install pandas openpyxl",
                "data": {},
            }

        try:
            # Try with pandas first for better data analysis
            if pandas_available:
                try:
                    # Read all sheets
                    excel_data = pd.read_excel(
                        file_path, sheet_name=None, engine="openpyxl"
                    )
                    sheets_analysis = {}
                    total_content = ""

                    for sheet_name, df in excel_data.items():
                        # Basic data analysis
                        sheet_analysis = {
                            "name": sheet_name,
                            "shape": df.shape,
                            "columns": list(df.columns),
                            "dtypes": df.dtypes.to_dict(),
                            "null_counts": df.isnull().sum().to_dict(),
                            "sample_data": df.head(3).to_dict() if not df.empty else {},
                        }

                        # Add basic statistics for numeric columns
                        numeric_cols = df.select_dtypes(include=["number"]).columns
                        if len(numeric_cols) > 0:
                            sheet_analysis["numeric_summary"] = (
                                df[numeric_cols].describe().to_dict()
                            )

                        sheets_analysis[sheet_name] = sheet_analysis

                        # Create text representation
                        sheet_text = f"\n=== Sheet: {sheet_name} ===\n"
                        sheet_text += (
                            f"Dimensions: {df.shape[0]} rows × {df.shape[1]} columns\n"
                        )
                        sheet_text += f"Columns: {', '.join(df.columns)}\n"

                        if not df.empty:
                            sheet_text += "\nFirst few rows:\n"
                            sheet_text += df.head(5).to_string() + "\n"

                        total_content += sheet_text

                    # Overall analysis
                    analysis = {
                        "total_sheets": len(excel_data),
                        "total_rows": sum(df.shape[0] for df in excel_data.values()),
                        "total_columns": sum(df.shape[1] for df in excel_data.values()),
                        "sheets_analysis": sheets_analysis,
                    }

                    content = f"""EXCEL WORKBOOK SUMMARY:
Total Sheets: {analysis["total_sheets"]}
Total Rows: {analysis["total_rows"]}
Total Columns: {analysis["total_columns"]}

WORKBOOK CONTENT:
{total_content}"""

                    return {
                        "success": True,
                        "data": {
                            "content_type": "spreadsheet",
                            "raw_content": content,
                            "file_path": str(file_path),
                            "file_size": file_path.stat().st_size,
                            "extraction_method": "excel_pandas",
                            "spreadsheet_analysis": analysis,
                        },
                    }

                except Exception as pandas_error:
                    # Fall back to openpyxl
                    pass

            # Fallback to openpyxl direct reading
            if openpyxl_available:
                workbook = load_workbook(file_path, data_only=True)
                sheets_content = []
                total_content = ""

                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    sheet_text = f"\n=== Sheet: {sheet_name} ===\n"

                    # Extract cell values (limit to reasonable size)
                    rows_content = []
                    max_rows = min(sheet.max_row, 100)  # Limit to first 100 rows
                    max_cols = min(sheet.max_column, 20)  # Limit to first 20 columns

                    for row in sheet.iter_rows(
                        min_row=1,
                        max_row=max_rows,
                        min_col=1,
                        max_col=max_cols,
                        values_only=True,
                    ):
                        row_values = [
                            str(cell) if cell is not None else "" for cell in row
                        ]
                        if any(row_values):  # Skip empty rows
                            rows_content.append(" | ".join(row_values))

                    sheet_text += "\n".join(rows_content[:20])  # First 20 rows
                    sheets_content.append(
                        {
                            "name": sheet_name,
                            "rows": sheet.max_row,
                            "columns": sheet.max_column,
                            "content_preview": sheet_text,
                        }
                    )
                    total_content += sheet_text + "\n"

                analysis = {
                    "total_sheets": len(workbook.sheetnames),
                    "sheets": sheets_content,
                    "extraction_method": "openpyxl_fallback",
                }

                content = f"""EXCEL WORKBOOK SUMMARY (Basic Extraction):
Total Sheets: {len(workbook.sheetnames)}
Sheet Names: {", ".join(workbook.sheetnames)}

WORKBOOK CONTENT:
{total_content}"""

                return {
                    "success": True,
                    "data": {
                        "content_type": "spreadsheet",
                        "raw_content": content,
                        "file_path": str(file_path),
                        "file_size": file_path.stat().st_size,
                        "extraction_method": "excel_openpyxl",
                        "spreadsheet_analysis": analysis,
                    },
                }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract Excel: {e}",
                "data": {},
            }

    def _extract_csv(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from CSV/TSV file (enhanced from mark-mate)"""
        try:
            # Detect delimiter
            delimiter = "\t" if file_path.suffix.lower() == ".tsv" else ","

            # Try different encodings
            encodings = [
                "utf-8",
                "utf-16",
                "latin-1",
                "cp1252",
                "cp1251",
                "gb2312",
                "shift_jis",
            ]

            df = None
            used_encoding = None

            # Try pandas first for better analysis
            if pandas_available:
                for encoding in encodings:
                    try:
                        df = pd.read_csv(
                            file_path, delimiter=delimiter, encoding=encoding
                        )
                        used_encoding = encoding
                        break
                    except UnicodeDecodeError:
                        continue
                    except Exception as e:
                        continue

            if df is not None:
                # Pandas-based analysis
                analysis = {
                    "shape": df.shape,
                    "columns": list(df.columns),
                    "dtypes": df.dtypes.to_dict(),
                    "null_counts": df.isnull().sum().to_dict(),
                    "null_percentage": (df.isnull().sum() / len(df) * 100).to_dict(),
                    "encoding_used": used_encoding,
                    "delimiter": "tab" if delimiter == "\t" else "comma",
                }

                # Add statistics for numeric columns
                numeric_cols = df.select_dtypes(include=["number"]).columns
                if len(numeric_cols) > 0:
                    analysis["numeric_summary"] = df[numeric_cols].describe().to_dict()

                # Add sample data
                analysis["sample_data"] = {
                    "first_5_rows": df.head(5).to_dict(),
                    "last_5_rows": df.tail(5).to_dict() if len(df) > 5 else {},
                }

                # Create content text
                content = f"""CSV/TSV DATA SUMMARY:
File: {file_path.name}
Dimensions: {df.shape[0]} rows × {df.shape[1]} columns
Encoding: {used_encoding}
Delimiter: {analysis["delimiter"]}

COLUMNS:
{", ".join(df.columns)}

DATA QUALITY:
"""

                # Add data quality information
                for col in df.columns:
                    null_percentage = analysis.get("null_percentage", {})
                    null_pct = null_percentage.get(col, 0)
                    content += f"- {col}: {null_pct:.1f}% missing\n"

                content += f"\nFIRST FEW ROWS:\n{df.head(10).to_string()}\n"

                if len(df) > 10:
                    content += f"\nLAST FEW ROWS:\n{df.tail(5).to_string()}\n"

                return {
                    "success": True,
                    "data": {
                        "content_type": "data",
                        "raw_content": content,
                        "file_path": str(file_path),
                        "file_size": file_path.stat().st_size,
                        "extraction_method": "csv_pandas",
                        "data_analysis": analysis,
                    },
                }
            else:
                # Fallback: basic CSV reading
                rows = []
                for encoding in encodings:
                    try:
                        with open(file_path, encoding=encoding) as f:
                            import csv

                            reader = csv.reader(f, delimiter=delimiter)
                            rows = list(reader)
                            used_encoding = encoding
                            break
                    except UnicodeDecodeError:
                        continue
                    except Exception as e:
                        continue

                if not rows:
                    return {
                        "success": False,
                        "error": "Could not read CSV file with any encoding",
                        "data": {},
                    }

                headers = rows[0] if rows else []
                data_rows = rows[1:] if len(rows) > 1 else []

                # Basic analysis
                analysis = {
                    "shape": (len(data_rows), len(headers)),
                    "columns": headers,
                    "total_rows": len(data_rows),
                    "encoding_used": used_encoding,
                    "delimiter": repr(delimiter),
                    "analysis_method": "basic_csv_reader",
                }

                # Create content text
                content = f"""CSV/TSV DATA SUMMARY:
File: {file_path.name}
Dimensions: {len(data_rows)} rows × {len(headers)} columns
Encoding: {used_encoding}
Delimiter: {repr(delimiter)}

COLUMNS:
{", ".join(headers)}

FIRST FEW ROWS:
"""
                # Show first 10 rows
                for i, row in enumerate(data_rows[:10]):
                    content += f"Row {i + 1}: {', '.join(row[:5])}{'...' if len(row) > 5 else ''}\n"

                return {
                    "success": True,
                    "data": {
                        "content_type": "data",
                        "raw_content": content,
                        "file_path": str(file_path),
                        "file_size": file_path.stat().st_size,
                        "extraction_method": "csv_basic",
                        "data_analysis": analysis,
                    },
                }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract CSV: {e}",
                "data": {},
            }

    def _extract_json(self, file_path: Path) -> Dict[str, Any]:
        """Extract text from JSON file"""
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                data = json.load(f)

            # Extract text recursively from JSON structure
            extracted_texts = []
            self._extract_text_from_json_recursive(data, extracted_texts)

            if not extracted_texts:
                return {
                    "success": False,
                    "error": "No text content found in JSON",
                    "data": {},
                }

            content = "\n\n".join(extracted_texts)
            return {
                "success": True,
                "data": {
                    "content_type": "text",
                    "raw_content": content,
                    "file_path": str(file_path),
                    "file_size": file_path.stat().st_size,
                    "extraction_method": "json",
                },
            }

        except json.JSONDecodeError as e:
            return {"success": False, "error": f"Invalid JSON format: {e}", "data": {}}
        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to process JSON: {e}",
                "data": {},
            }

    def _extract_text_from_json_recursive(self, obj, texts):
        """Recursively extract text from JSON object"""
        if isinstance(obj, dict):
            for value in obj.values():
                if isinstance(value, str) and len(value.strip()) > 10:
                    texts.append(value.strip())
                elif isinstance(value, (dict, list)):
                    self._extract_text_from_json_recursive(value, texts)
        elif isinstance(obj, list):
            for item in obj:
                if isinstance(item, str) and len(item.strip()) > 10:
                    texts.append(item.strip())
                elif isinstance(item, (dict, list)):
                    self._extract_text_from_json_recursive(item, texts)
