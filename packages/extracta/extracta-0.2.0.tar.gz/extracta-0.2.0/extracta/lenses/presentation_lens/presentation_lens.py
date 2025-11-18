"""
Presentation lens for analyzing presentation files (PPTX, etc.)

This lens specializes in presentation formats, extracting both textual content
and visual elements for comprehensive analysis.
"""

from pathlib import Path
from typing import Dict, Any, List

from ..base_lens import BaseLens

# Import optional libraries for presentation processing
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import fitz  # PyMuPDF for PDF rendering
except ImportError:
    fitz = None


class PresentationLens(BaseLens):
    """Lens for extracting content from presentation files (slides, videos, demos)

    This lens handles both traditional slide presentations (PPTX) and video presentations
    (screencasts, recorded demos). It provides semantic flexibility for users who think
    of presentations broadly, while maintaining technical routing internally.
    """

    SUPPORTED_EXTENSIONS = {
        ".pptx",  # PowerPoint presentations
        ".ppt",  # Legacy PowerPoint
        ".mp4",  # Video presentations/screencasts
        ".mov",  # QuickTime videos
        ".avi",  # AVI videos
        ".webm",  # WebM videos
        ".mkv",  # Matroska videos
        ".flv",  # Flash videos
        ".wmv",  # Windows Media videos
        ".m4v",  # MPEG-4 videos
        # Future: .odp (OpenDocument), .key (Keynote), etc.
    }

    def __init__(self, extract_images: bool = False, render_slides: bool = False):
        """Initialize presentation lens.

        Args:
            extract_images: Whether to extract images from slides
            render_slides: Whether to render slides as images for analysis
        """
        self.extract_images = extract_images
        self.render_slides = render_slides
        self.pptx_available = Presentation is not None
        self.pdf_available = fitz is not None

    def extract(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from presentation file (slides or video)"""
        try:
            # Ensure file_path is a Path object
            if isinstance(file_path, str):
                file_path = Path(file_path)

            if file_path.suffix.lower() not in self.SUPPORTED_EXTENSIONS:
                return {
                    "success": False,
                    "error": f"Unsupported presentation format: {file_path.suffix}",
                    "data": {},
                }

            # Route based on content type
            video_extensions = {
                ".mp4",
                ".mov",
                ".avi",
                ".webm",
                ".mkv",
                ".flv",
                ".wmv",
                ".m4v",
            }

            if file_path.suffix.lower() in video_extensions:
                # Route video presentations to video lens
                return self._extract_video_presentation(file_path)
            elif file_path.suffix.lower() in [".pptx", ".ppt"]:
                return self._extract_pptx(file_path)
            else:
                return {
                    "success": False,
                    "error": f"Unsupported presentation format: {file_path.suffix}",
                    "data": {},
                }

        except Exception as e:
            return {"success": False, "error": str(e), "data": {}}

    def _extract_video_presentation(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from video presentation by routing to video lens"""
        try:
            # Import video lens and delegate
            from .video_lens import VideoLens

            video_lens = VideoLens()

            # Extract video content
            result = video_lens.extract(file_path)

            if result["success"]:
                # Enhance the result to indicate this was processed as a presentation
                result["data"]["presentation_context"] = {
                    "semantic_type": "video_presentation",
                    "routed_from": "presentation_lens",
                    "description": "Video content processed as presentation (screencast/demo)",
                }

                # Add presentation-specific analysis if applicable
                if "transcript" in result["data"]:
                    # Analyze presentation structure from transcript
                    presentation_analysis = self._analyze_video_presentation(
                        result["data"]
                    )
                    result["data"]["presentation_analysis"] = presentation_analysis

            return result

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract video presentation: {e}",
                "data": {},
            }

    def _extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extract content from PowerPoint presentation"""
        if not self.pptx_available:
            return {
                "success": False,
                "error": "PPTX processing not available. Install with: pip install python-pptx",
                "data": {},
            }

        try:
            prs = Presentation(file_path)

            # Extract basic presentation metadata
            metadata = self._extract_presentation_metadata(prs, file_path)

            # Extract slide content
            slides_content = self._extract_slides_content(prs)

            # Separate slide text from presenter notes
            slide_text = []
            presenter_notes = []

            for slide_data in slides_content:
                if slide_data.get("slide_text"):
                    slide_text.append(slide_data["slide_text"])
                if slide_data.get("notes_text"):
                    presenter_notes.append(slide_data["notes_text"])

            # Combine all slide text
            combined_slide_text = "\n\n".join(slide_text)
            combined_notes_text = "\n\n".join(presenter_notes)

            # Extract images if requested
            extracted_images = []
            if self.extract_images:
                extracted_images = self._extract_slide_images(prs, file_path)

            # Render slides as images if requested
            rendered_slides = []
            if self.render_slides:
                rendered_slides = self._render_slides_to_images(prs, file_path)

            data = {
                "content_type": "presentation",
                "format": "pptx",
                "file_path": str(file_path),
                "file_size": file_path.stat().st_size,
                "metadata": metadata,
                "slide_text": combined_slide_text,
                "presenter_notes": combined_notes_text,
                "slides_content": slides_content,
                "slide_count": len(slides_content),
                "extracted_images": extracted_images,
                "rendered_slides": rendered_slides,
                "extraction_method": "pptx",
            }

            # Apply cascading analysis for code detection
            from .. import analyze_extracted_content

            enhanced_data = analyze_extracted_content(data)
            if "cascading_analysis" in enhanced_data:
                data["cascading_analysis"] = enhanced_data["cascading_analysis"]

            return {
                "success": True,
                "data": data,
            }

        except Exception as e:
            return {
                "success": False,
                "error": f"Failed to extract PPTX: {e}",
                "data": {},
            }

    def _extract_presentation_metadata(
        self, prs: Presentation, file_path: Path
    ) -> Dict[str, Any]:
        """Extract presentation metadata"""
        metadata = {
            "slide_count": len(prs.slides),
            "format": "pptx",
        }

        # Extract core properties if available
        if hasattr(prs, "core_properties"):
            props = prs.core_properties
            metadata.update(
                {
                    "title": str(props.title or ""),
                    "author": str(props.author or ""),
                    "subject": str(props.subject or ""),
                    "keywords": str(props.keywords or ""),
                    "comments": str(props.comments or ""),
                    "created": props.created.isoformat() if props.created else None,
                    "modified": props.modified.isoformat() if props.modified else None,
                }
            )

        return metadata

    def _analyze_video_presentation(self, video_data: Dict[str, Any]) -> Dict[str, Any]:
        """Analyze video content in presentation context"""
        analysis = {
            "presentation_type": "video_presentation",
            "content_structure": {},
            "engagement_indicators": {},
        }

        # Analyze transcript for presentation structure
        if "transcript" in video_data and video_data["transcript"]:
            transcript = video_data["transcript"]

            # Basic presentation structure analysis
            sentences = [s.strip() for s in transcript.split(".") if s.strip()]
            analysis["content_structure"] = {
                "sentence_count": len(sentences),
                "estimated_duration": video_data.get("duration", 0),
                "content_density": len(sentences)
                / max(video_data.get("duration", 1), 1),
            }

            # Look for presentation indicators in transcript
            presentation_keywords = [
                "welcome",
                "introduction",
                "overview",
                "summary",
                "conclusion",
                "next",
                "firstly",
                "secondly",
                "finally",
                "thank you",
                "questions",
                "q&a",
                "demo",
                "example",
                "case study",
            ]

            keyword_count = sum(
                1
                for keyword in presentation_keywords
                if keyword.lower() in transcript.lower()
            )

            analysis["engagement_indicators"] = {
                "presentation_keywords_found": keyword_count,
                "structure_indicators": min(keyword_count, 5),  # Cap at 5 for scoring
            }

        return analysis

    def _extract_slides_content(self, prs: Presentation) -> List[Dict[str, Any]]:
        """Extract content from each slide"""
        slides_content = []

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "slide_number": i + 1,
                "slide_text": "",
                "notes_text": "",
                "shape_count": len(slide.shapes),
                "has_title": False,
                "has_content": False,
            }

            # Extract text from shapes
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

                    # Try to identify title vs content (basic heuristic)
                    if not slide_data["has_title"] and len(shape.text.strip()) < 100:
                        slide_data["has_title"] = True
                    else:
                        slide_data["has_content"] = True

            slide_data["slide_text"] = "\n".join(text_parts)

            # Extract presenter notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text_parts = []
                for shape in notes_slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        notes_text_parts.append(shape.text.strip())
                slide_data["notes_text"] = "\n".join(notes_text_parts)

            slides_content.append(slide_data)

        return slides_content

    def _extract_slide_images(
        self, prs: Presentation, file_path: Path
    ) -> List[Dict[str, Any]]:
        """Extract images from slides"""
        # This is a placeholder for image extraction
        # In a full implementation, this would extract embedded images
        # and save them to temporary files for analysis
        return []

    def _render_slides_to_images(
        self, prs: Presentation, file_path: Path
    ) -> List[Dict[str, Any]]:
        """Render slides as images for visual analysis"""
        # This is a placeholder for slide rendering
        # In a full implementation, this would:
        # 1. Use python-pptx to render slides to images
        # 2. Save images to temporary files
        # 3. Return paths for image analysis
        return []

    def analyze_presentation_structure(
        self, slides_content: List[Dict[str, Any]]
    ) -> Dict[str, Any]:
        """Analyze the structure and flow of the presentation"""
        if not slides_content:
            return {}

        # Analyze title slide patterns
        title_slides = [
            s for s in slides_content if s.get("has_title") and not s.get("has_content")
        ]

        # Analyze content distribution
        slides_with_content = [s for s in slides_content if s.get("slide_text")]
        slides_with_notes = [s for s in slides_content if s.get("notes_text")]

        # Calculate average text length per slide
        text_lengths = [len(s.get("slide_text", "")) for s in slides_content]
        avg_text_length = sum(text_lengths) / len(text_lengths) if text_lengths else 0

        return {
            "title_slides": len(title_slides),
            "content_slides": len(slides_with_content),
            "slides_with_notes": len(slides_with_notes),
            "avg_text_per_slide": round(avg_text_length, 1),
            "notes_coverage": round(
                len(slides_with_notes) / len(slides_content) * 100, 1
            )
            if slides_content
            else 0,
        }
