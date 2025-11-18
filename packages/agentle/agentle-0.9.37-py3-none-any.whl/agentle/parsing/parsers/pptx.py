"""PowerPoint Presentation Parser Module

This module provides functionality for parsing Microsoft PowerPoint presentations
(.ppt, .pptx, .pptm) into structured representations. It can extract text content,
process embedded images, and organize the presentation by slides.

Enhancements implemented:
- File existence & extension validation
- Structured exception hierarchy
- Only legacy binary ``.ppt`` conversion (not converting ``.pptm`` to preserve macros)
- Safe temporary conversion without leaked files
- Group shape recursion (images/text in grouped shapes)
- Optional speaker notes extraction
- Optional whitespace normalization
- Robust image MIME type detection (uses real content type; fallback sniffing)
- Handling of shapes without extensions (no empty suffix issues)
- Concurrency & retry with timeout for image description calls
- Caching of duplicate images
- Progress callback & metrics collection
- Logging & warnings when strategy is high but provider missing
- Graceful continuation on provider errors (configurable)
"""

import os
import subprocess
import tempfile
import asyncio
import logging
import time
import atexit
from dataclasses import dataclass
from collections.abc import MutableSequence
from pathlib import Path
from typing import (
    Callable,
    Literal,
    cast,
    Any,
    Iterable,
    Awaitable,
    Tuple,
    List,
)
from io import BytesIO

try:  # imghdr may be deprecated in future versions; treat as optional
    import imghdr  # type: ignore
except Exception:  # pragma: no cover
    imghdr = None  # type: ignore

from rsb.models.field import Field

from agentle.generations.models.message_parts.file import FilePart
from agentle.generations.models.structured_outputs_store.visual_media_description import (
    VisualMediaDescription,
)
from agentle.generations.providers.base.generation_provider import GenerationProvider
from agentle.parsing.document_parser import DocumentParser
from agentle.parsing.image import Image
from agentle.parsing.parsed_file import ParsedFile
from agentle.parsing.section_content import SectionContent

logger = logging.getLogger(__name__)


# -----------------------------
# Exception Hierarchy
# -----------------------------
class PptxParserError(Exception):
    """Base error for PPTX parsing operations."""


class FileValidationError(PptxParserError):
    """Raised when the input file is missing or unsupported."""


class DependencyMissingError(PptxParserError):
    """Raised when an external dependency (e.g., LibreOffice, python-pptx) is missing."""


class ConversionError(PptxParserError):
    """Raised when legacy PPT conversion fails."""


class CorruptFileError(PptxParserError):
    """Raised when the PPTX container appears invalid or corrupted."""


class ProviderError(PptxParserError):
    """Raised when the visual description provider fails (fatal)."""


@dataclass
class ImageMetrics:
    slide_index: int
    image_index: int
    bytes: int
    duration_s: float
    retries: int
    success: bool
    error: str | None = None


@dataclass
class ParseMetrics:
    total_slides: int = 0
    images_processed: int = 0
    images_described: int = 0
    start_time_s: float = 0.0
    end_time_s: float = 0.0
    image_metrics: list[ImageMetrics] | None = None

    @property
    def total_duration_s(self) -> float:
        return self.end_time_s - self.start_time_s if self.end_time_s else 0.0


class PptxFileParser(DocumentParser):
    """
    Parser for processing Microsoft PowerPoint presentations (.ppt, .pptx, .pptm).

    This parser extracts content from PowerPoint presentations, including text and embedded
    images. Each slide in the presentation is parsed as a separate section in the resulting
    ParsedFile. With the "high" strategy, embedded images are analyzed using a visual
    description agent to extract text via OCR and generate descriptions.

    The parser supports both modern PowerPoint formats (.pptx, .pptm) and the legacy
    format (.ppt), converting the latter to .pptx using LibreOffice when necessary.

    **Attributes:**

    *   `strategy` (Literal["high", "low"]):
        The parsing strategy to use. Defaults to "high".
        - "high": Performs thorough parsing including OCR and image analysis
        - "low": Performs basic text extraction without analyzing images

        **Example:**
        ```python
        parser = PptxFileParser(strategy="low")  # Use faster, less intensive parsing
        ```

    *   `visual_description_agent` (Agent[VisualMediaDescription]):
        The agent used to analyze and describe image content. If provided and
        strategy is "high", this agent will be used to analyze images embedded
        in the presentation.
        Defaults to the agent created by `visual_description_agent_default_factory()`.

        **Example:**
        ```python
        from agentle.agents.agent import Agent
        from agentle.generations.models.structured_outputs_store.visual_media_description import VisualMediaDescription

        custom_agent = Agent(
            model="gemini-2.0-pro-vision",
            instructions="Focus on diagram and chart analysis in presentations",
            response_schema=VisualMediaDescription
        )

        parser = PptxFileParser(visual_description_agent=custom_agent)
        ```

    *   `multi_modal_provider` (GenerationProvider):
        An alternative to using a visual_description_agent. This is a generation
        provider capable of handling multi-modal content (text and images).
        Defaults to GoogleGenerationProvider().

        Note: You cannot use both visual_description_agent and multi_modal_provider
        at the same time.

    **Usage Examples:**

    Basic parsing of a PowerPoint presentation:
    ```python
    from agentle.parsing.parsers.pptx import PptxFileParser

    # Create a parser with default settings
    parser = PptxFileParser()

    # Parse a PowerPoint file
    parsed_presentation = parser.parse("presentation.pptx")

    # Access the slides (as sections)
    for i, section in enumerate(parsed_presentation.sections):
        print(f"Slide {i+1} content:")
        print(section.text[:100] + "...")  # Print first 100 chars of each slide
    ```

    Processing images in a presentation:
    ```python
    from agentle.parsing.parsers.pptx import PptxFileParser

    # Create a parser with high-detail strategy
    parser = PptxFileParser(strategy="high")

    # Parse a presentation with images
    slide_deck = parser.parse("slide_deck.pptx")

    # Extract and process images
    for i, section in enumerate(slide_deck.sections):
        slide_num = i + 1
        print(f"Slide {slide_num} has {len(section.images)} images")

        for j, image in enumerate(section.images):
            print(f"  Image {j+1}:")
            if image.ocr_text:
                print(f"    OCR text: {image.ocr_text}")
    ```

    **Requirements:**

    For parsing .ppt (legacy format) files, LibreOffice must be installed on the system
    to perform the conversion to .pptx. If LibreOffice is not installed, a RuntimeError
    will be raised when attempting to parse .ppt files.
    """

    # Backward compatibility & clearer naming
    parser_type: Literal["pptx"] = "pptx"
    type: Literal["pptx"] = "pptx"  # alias (deprecated)

    strategy: Literal["high", "low"] = Field(default="high")

    visual_description_provider: GenerationProvider | None = Field(
        default=None,
    )
    """
    The agent to use for generating the visual description of the document.
    Useful when you want to customize the prompt for the visual description.
    """

    # Additional configuration fields
    include_notes: bool = Field(default=False)
    normalize_whitespace: bool = Field(default=True)
    continue_on_provider_error: bool = Field(default=True)
    image_description_retries: int = Field(default=2)
    image_description_timeout: float = Field(default=25.0)
    max_concurrent_image_tasks: int = Field(default=4)
    max_image_bytes: int | None = Field(default=None)
    progress_callback: Callable[[int, int], None] | None = Field(default=None)
    collect_metrics: bool = Field(default=True)

    # Internal last metrics (publicly readable after parse)
    last_parse_metrics: ParseMetrics | None = None

    async def parse_async(self, document_path: str) -> ParsedFile:
        """
        Asynchronously parse a PowerPoint presentation and generate a structured representation.

        This method reads a PowerPoint file, extracts text and image content from each slide,
        and processes embedded images to extract text and generate descriptions when using
        the "high" strategy.

        For .ppt (legacy format) files, the method first converts them to .pptx format
        using LibreOffice before processing.

        Args:
            document_path (str): Path to the PowerPoint file to be parsed

        Returns:
            ParsedFile: A structured representation where:
                - Each slide is a separate section
                - Text content is extracted from each slide
                - Images are extracted and (optionally) analyzed

        Raises:
            RuntimeError: If converting a .ppt file fails (e.g., if LibreOffice is not installed)
                or if the conversion times out

        Example:
            ```python
            import asyncio
            from agentle.parsing.parsers.pptx import PptxFileParser

            async def process_presentation():
                parser = PptxFileParser(strategy="high")
                result = await parser.parse_async("slide_deck.pptx")

                # Print information about the slides
                print(f"Presentation has {len(result.sections)} slides")

                # Access slide content
                for i, section in enumerate(result.sections):
                    print(f"Slide {i+1} content:")
                    print(section.text[:100] + "...")

                    # Check for images
                    if section.images:
                        print(f"  Contains {len(section.images)} images")

            asyncio.run(process_presentation())
            ```

        Note:
            This method uses the python-pptx library to read PowerPoint files. For optimal
            results with legacy .ppt files, ensure LibreOffice is installed on the system.
        """
        import hashlib

        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.presentation import Presentation as PptxPresentation
            from pptx.shapes.autoshape import Shape
            from pptx.shapes.picture import Picture
        except ModuleNotFoundError as e:  # pragma: no cover - environment dependent
            raise DependencyMissingError(
                "python-pptx is required. Install with: pip install python-pptx"
            ) from e

        path = Path(document_path)
        if not path.exists() or not path.is_file():
            raise FileValidationError(f"File not found: {document_path}")

        suffix = path.suffix.lower()
        supported = {".ppt", ".pptx", ".pptm"}
        if suffix not in supported:
            raise FileValidationError(
                f"Unsupported extension '{suffix}'. Supported: {', '.join(sorted(supported))}"
            )

        # Convert only legacy binary .ppt
        prs_data: bytes | None = None
        if suffix == ".ppt":
            logger.debug(
                "Converting legacy .ppt file to .pptx bytes: %s", document_path
            )
            prs_data = self._convert_legacy_ppt_to_pptx_bytes(document_path)
        elif suffix == ".pptm":
            logger.debug(
                "Opening .pptm directly without conversion (macros ignored): %s",
                document_path,
            )

        try:
            if prs_data is not None:
                prs: PptxPresentation = Presentation(BytesIO(prs_data))
            else:
                prs = Presentation(document_path)
        except Exception as e:  # python-pptx raises diverse errors
            raise CorruptFileError(
                f"Failed to open PowerPoint file '{document_path}': {e}"
            ) from e

        total_slides = len(prs.slides)
        if total_slides == 0:
            logger.warning("Presentation contains zero slides: %s", document_path)

        if self.strategy == "high" and not self.visual_description_provider:
            logger.warning(
                "High strategy requested but no visual_description_provider configured; skipping image analysis."
            )

        metrics = ParseMetrics(
            start_time_s=time.time(),
            total_slides=total_slides,
            image_metrics=[] if self.collect_metrics else None,
        )
        sections: MutableSequence[SectionContent] = []
        processed_images: dict[str, tuple[str, str]] = {}  # hash -> (md, ocr)

        # Helper for progress callback
        def report_progress(idx: int):
            if self.progress_callback:
                try:
                    self.progress_callback(idx, total_slides)
                except Exception:  # pragma: no cover
                    logger.debug("Progress callback raised an exception", exc_info=True)

        for slide_index, slide in enumerate(prs.slides, start=1):
            report_progress(slide_index)
            slide_texts: list[str] = []
            slide_images: list[
                tuple[str, bytes, str, str | None]
            ] = []  # (name, data, hash, content_type)

            # Recursive traversal for grouped shapes
            def iter_shapes(
                shapes: Any,
            ) -> Iterable[Any]:  # shapes is a pptx collection
                try:
                    for shp in shapes:  # type: ignore[assignment]
                        yield shp
                        nested = getattr(shp, "shapes", None)
                        if (
                            nested is not None
                            and nested is not shapes
                            and nested is not shp
                        ):  # group/container
                            yield from iter_shapes(nested)
                except TypeError:  # not iterable
                    return

            for shape in iter_shapes(slide.shapes):
                try:
                    if getattr(shape, "has_text_frame", False):
                        shape_with_text = cast(Shape, shape)
                        text_str: str = shape_with_text.text
                        if text_str:
                            slide_texts.append(text_str)

                    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                        picture_shape = cast(Picture, shape)
                        image_blob: bytes = picture_shape.image.blob
                        if (
                            self.max_image_bytes
                            and len(image_blob) > self.max_image_bytes
                        ):
                            logger.info(
                                "Skipping large image (>%d bytes) on slide %d",
                                self.max_image_bytes,
                                slide_index,
                            )
                            continue
                        image_hash = hashlib.sha256(image_blob).hexdigest()
                        raw_name = (
                            getattr(shape, "name", None)
                            or f"slide_{slide_index}_img_{image_hash[:8]}"
                        )
                        content_type = getattr(
                            picture_shape.image, "content_type", None
                        )
                        slide_images.append(
                            (raw_name, image_blob, image_hash, content_type)
                        )
                except Exception:  # robust against unexpected shape edge cases
                    logger.debug(
                        "Error processing shape on slide %d", slide_index, exc_info=True
                    )

            # Speaker notes (optional)
            if self.include_notes:
                try:
                    if (
                        getattr(slide, "has_notes_slide", False)
                        and slide.notes_slide
                        and slide.notes_slide.notes_text_frame
                    ):  # type: ignore[attr-defined]
                        notes_text = slide.notes_slide.notes_text_frame.text  # type: ignore[attr-defined]
                        if notes_text:
                            slide_texts.append("--- Notes ---\n" + notes_text)
                except Exception:  # pragma: no cover
                    logger.debug(
                        "Failed to extract notes for slide %d",
                        slide_index,
                        exc_info=True,
                    )

            combined_text: str = "\n".join(slide_texts)
            final_images: MutableSequence[Image] = []
            image_descriptions: MutableSequence[str] = []
            provider = (
                self.visual_description_provider if self.strategy == "high" else None
            )

            # Image description logic with concurrency, retry & timeout
            if provider and slide_images:
                semaphore = asyncio.Semaphore(self.max_concurrent_image_tasks)
                describe_tasks: List[
                    Awaitable[Tuple[int, bool, str, str, Image, str | None]]
                ] = []

                async def describe_one(
                    img_idx: int,
                    image_name: str,
                    image_blob: bytes,
                    image_hash: str,
                    content_type: str | None,
                ) -> Tuple[int, bool, str, str, Image, str | None]:
                    nonlocal metrics
                    start = time.time()
                    retries_used = 0
                    if image_hash in processed_images:
                        cached_md, cached_ocr = processed_images[image_hash]
                        if metrics.image_metrics is not None:
                            metrics.image_metrics.append(
                                ImageMetrics(
                                    slide_index,
                                    img_idx,
                                    len(image_blob),
                                    0.0,
                                    0,
                                    True,
                                    None,
                                )
                            )
                        return (
                            img_idx,
                            True,
                            cached_md,
                            cached_ocr,
                            Image(
                                name=image_name,
                                contents=image_blob,
                                ocr_text=cached_ocr,
                            ),
                            None,
                        )
                    mime_type = self._deduce_mime_type(
                        image_name, image_blob, content_type
                    )
                    file_part = FilePart(mime_type=mime_type, data=image_blob)
                    attempt_error: Exception | None = None
                    for attempt in range(self.image_description_retries + 1):
                        retries_used = attempt
                        try:
                            async with semaphore:
                                assert provider is not None  # for type checker
                                response = await asyncio.wait_for(
                                    provider.generate_by_prompt_async(  # type: ignore[attr-defined]
                                        file_part,
                                        developer_prompt="You are a helpful assistant that deeply understands visual media.",
                                        response_schema=VisualMediaDescription,
                                    ),
                                    timeout=self.image_description_timeout,
                                )
                            image_md: str = response.parsed.md
                            image_ocr = response.parsed.ocr_text
                            processed_images[image_hash] = (image_md, image_ocr or "")
                            duration = time.time() - start
                            if metrics.image_metrics is not None:
                                metrics.image_metrics.append(
                                    ImageMetrics(
                                        slide_index,
                                        img_idx,
                                        len(image_blob),
                                        duration,
                                        retries_used,
                                        True,
                                        None,
                                    )
                                )
                            return (
                                img_idx,
                                True,
                                image_md,
                                image_ocr or "",
                                Image(
                                    name=image_name,
                                    contents=image_blob,
                                    ocr_text=image_ocr,
                                ),
                                None,
                            )
                        except (
                            Exception
                        ) as exc:  # pragma: no cover - diverse provider errors
                            attempt_error = exc
                            await asyncio.sleep(0.2 * (attempt + 1))
                    duration = time.time() - start
                    err_str = (
                        repr(attempt_error)
                        if attempt_error
                        else "Unknown provider failure"
                    )
                    if metrics.image_metrics is not None:
                        metrics.image_metrics.append(
                            ImageMetrics(
                                slide_index,
                                img_idx,
                                len(image_blob),
                                duration,
                                retries_used,
                                False,
                                err_str,
                            )
                        )
                    if not self.continue_on_provider_error:
                        raise ProviderError(
                            f"Image description failed: {err_str}"
                        ) from attempt_error
                    logger.warning(
                        "Continuing after provider failure on slide %d image %d: %s",
                        slide_index,
                        img_idx,
                        err_str,
                    )
                    return (
                        img_idx,
                        False,
                        "(Description unavailable due to error)",
                        "",
                        Image(name=image_name, contents=image_blob, ocr_text=None),
                        err_str,
                    )

                for img_idx, (
                    image_name,
                    image_blob,
                    image_hash,
                    content_type,
                ) in enumerate(slide_images, start=1):
                    describe_tasks.append(
                        describe_one(
                            img_idx, image_name, image_blob, image_hash, content_type
                        )
                    )

                results: List[
                    Tuple[int, bool, str, str, Image, str | None]
                ] = await asyncio.gather(*describe_tasks, return_exceptions=False)
                results.sort(key=lambda r: r[0])
                for img_idx, _success, md, _ocr, image_obj, _err in results:
                    image_descriptions.append(
                        f"Slide {slide_index} - Image {img_idx}: {md}"
                    )
                    final_images.append(image_obj)
                if metrics.image_metrics is not None:
                    metrics.images_processed += len(results)
                    metrics.images_described += sum(1 for r in results if r[1])

            if image_descriptions:
                combined_text += "\n\n" + "\n".join(image_descriptions)

            if self.normalize_whitespace:
                combined_text = self._normalize_whitespace(combined_text)

            section_content = SectionContent(
                number=slide_index,
                text=combined_text,
                md=combined_text,
                images=final_images,
            )
            sections.append(section_content)

        metrics.end_time_s = time.time()
        self.last_parse_metrics = metrics if self.collect_metrics else None

        return ParsedFile(name=path.name, sections=sections)

    # -------------------------------------------------------------
    # Helper Methods
    # -------------------------------------------------------------
    def _convert_legacy_ppt_to_pptx_bytes(self, document_path: str) -> bytes:
        """Convert legacy binary .ppt file to .pptx bytes using LibreOffice.

        Returns the raw bytes of the converted file (in-memory) so no temp file leakage occurs.
        """

        def _is_libreoffice_installed() -> bool:
            try:
                subprocess.run(
                    ["libreoffice", "--version"],
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                    check=True,
                )
                return True
            except (subprocess.CalledProcessError, FileNotFoundError):
                return False

        if not _is_libreoffice_installed():
            raise DependencyMissingError(
                "LibreOffice not found in system PATH for .ppt conversion"
            )

        with tempfile.TemporaryDirectory() as temp_dir:
            input_filename = Path(document_path).name
            input_path = os.path.join(temp_dir, input_filename)
            with (
                open(document_path, "rb") as src_file,
                open(input_path, "wb") as dst_file,
            ):
                dst_file.write(src_file.read())
            try:
                proc = subprocess.run(
                    [
                        "libreoffice",
                        "--headless",
                        "--convert-to",
                        "pptx",
                        "--outdir",
                        temp_dir,
                        input_path,
                    ],
                    check=True,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.PIPE,
                    timeout=90,
                )
            except subprocess.CalledProcessError as e:
                err = (
                    e.stderr.decode(errors="ignore").strip()
                    if e.stderr
                    else "Unknown error"
                )
                raise ConversionError(f"LibreOffice conversion failed: {err}") from e
            except subprocess.TimeoutExpired as e:
                raise ConversionError(
                    "LibreOffice conversion timed out after 90 seconds"
                ) from e

            output_filename = Path(input_filename).stem + ".pptx"
            output_path = os.path.join(temp_dir, output_filename)
            if not os.path.exists(output_path):
                files = os.listdir(temp_dir)
                raise ConversionError(
                    f"Converted file not found at {output_path}. Found: {files}. LibreOffice stdout: {proc.stdout[:200]!r}"
                )
            with open(output_path, "rb") as f:
                return f.read()

    def _deduce_mime_type(
        self, image_name: str, image_blob: bytes, declared: str | None
    ) -> str:
        """Best effort derivation of image MIME type.

        Order:
            1. Declared content_type from pptx
            2. imghdr sniff
            3. Extension mapping
            4. application/octet-stream fallback
        """
        if declared:
            return declared
        # imghdr
        if imghdr is not None:  # type: ignore[truthy-function]
            try:
                kind = imghdr.what(None, h=image_blob)  # type: ignore[attr-defined]
            except Exception:  # pragma: no cover
                kind = None
            if kind:
                return f"image/{kind.lower()}"
        ext = Path(image_name).suffix.lower()
        ext_map = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".bmp": "image/bmp",
            ".tiff": "image/tiff",
            ".webp": "image/webp",
        }
        if ext in ext_map:
            return ext_map[ext]
        return "application/octet-stream"

    def _normalize_whitespace(self, text: str) -> str:
        lines = [ln.rstrip() for ln in text.splitlines()]
        # Collapse multiple blank lines
        collapsed: list[str] = []
        blank = False
        for ln in lines:
            if ln.strip():
                collapsed.append(ln)
                blank = False
            else:
                if not blank:
                    collapsed.append("")
                blank = True
        return "\n".join(collapsed).strip()

    # Utility: optional cleanup registration for future expansions
    @staticmethod
    def _register_temp_file_for_cleanup(
        path: str,
    ) -> None:  # pragma: no cover (defensive)
        def _cleanup(p: str = path) -> None:
            try:
                if os.path.exists(p):
                    os.unlink(p)
            except Exception:
                logger.debug("Failed to cleanup temp file: %s", p, exc_info=True)

        atexit.register(_cleanup)
