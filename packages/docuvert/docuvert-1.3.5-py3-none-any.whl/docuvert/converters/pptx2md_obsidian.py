"""
Enhanced PPTX to Markdown converter optimized for Obsidian.
Preserves formatting, extracts images, and uses Obsidian-specific features.
"""

import sys
import os
import re
import base64
from pathlib import Path
from typing import Optional, List, Dict, Any, Tuple
from datetime import datetime

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from core.ast import ASTNode, NodeType, StyleInfo, DocumentMetadata, create_document, create_paragraph, create_heading
from core.exceptions import ConversionError, DependencyError


class Pptx2MdObsidianConverter:
    """
    Converts PPTX files to Markdown format optimized for Obsidian.
    
    Features:
    - Preserves slide structure and hierarchy
    - Extracts and embeds images
    - Maintains formatting (bold, italic, underline, colors)
    - Creates internal links between slides
    - Adds metadata as YAML frontmatter
    - Uses Obsidian callouts for special content
    - Preserves speaker notes as collapsible sections
    - Extracts charts and diagrams with descriptions
    - Maintains table formatting with Obsidian extensions
    - Creates slide thumbnails when possible
    """

    def __init__(self):
        """Initialize the converter with default settings."""
        self.image_dir = None
        self.image_counter = 0
        self.slide_titles = {}
        self.extracted_images = {}
        
    def parse_pptx2ast(self, input_path: str) -> ASTNode:
        """Parse PPTX file and convert to AST representation with enhanced features."""
        
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError as e:
            raise DependencyError(
                "python-pptx is required for PPTX conversion",
                missing_dependency="python-pptx"
            ) from e

        # Load presentation
        prs = Presentation(input_path)
        doc = create_document()
        
        # Set up image extraction directory
        self.image_dir = Path(input_path).parent / f"{Path(input_path).stem}_assets"
        print(f"\n🖼️  Image extraction directory: {self.image_dir}")
        
        # Extract comprehensive metadata
        metadata = self._extract_metadata(prs, input_path)
        doc.metadata = metadata

        # First pass: collect all slide titles for cross-referencing
        self._collect_slide_titles(prs)
        
        # Process slides with enhanced features
        for slide_idx, slide in enumerate(prs.slides):
            slide_ast = self._process_slide_enhanced(slide, slide_idx, prs)
            if slide_ast:
                doc.add_child(slide_ast)
        
        # Add presentation summary at the end
        summary = self._create_presentation_summary(prs)
        if summary:
            doc.add_child(summary)
            
        return doc

    def _extract_metadata(self, prs, input_path: str) -> DocumentMetadata:
        """Extract comprehensive metadata from presentation."""
        metadata = DocumentMetadata()
        
        if prs.core_properties:
            metadata.title = prs.core_properties.title or Path(input_path).stem
            metadata.author = prs.core_properties.author or "Unknown"
            metadata.subject = prs.core_properties.subject
            metadata.keywords = prs.core_properties.keywords
            metadata.created = prs.core_properties.created
            metadata.modified = prs.core_properties.modified
            metadata.category = prs.core_properties.category
            
        # Add custom Obsidian metadata
        metadata.custom = {
            "total_slides": len(prs.slides),
            "presentation_format": "PowerPoint",
            "conversion_date": datetime.now().isoformat(),
            "tags": ["presentation", "slides", "converted-from-pptx"],
            "aliases": [Path(input_path).stem, metadata.title] if metadata.title else [Path(input_path).stem]
        }
        
        return metadata

    def _collect_slide_titles(self, prs):
        """Collect all slide titles for cross-referencing."""
        for idx, slide in enumerate(prs.slides):
            title = self._extract_slide_title(slide)
            if title:
                self.slide_titles[idx] = title

    def _extract_slide_title(self, slide) -> Optional[str]:
        """Extract the title from a slide."""
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
        
        # Fallback: look for the first text in a title placeholder
        for shape in slide.shapes:
            if shape.is_placeholder:
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:  # Title
                    if hasattr(shape, 'text') and shape.text:
                        return shape.text.strip()
        return None

    def _process_slide_enhanced(self, slide, slide_idx: int, prs) -> ASTNode:
        """Process a slide with enhanced Obsidian features."""
        slide_container = ASTNode(NodeType.SECTION)
        
        # Add slide navigation header
        nav_header = self._create_slide_navigation(slide_idx, len(prs.slides))
        if nav_header:
            slide_container.add_child(nav_header)
        
        # Add slide title as heading
        title = self._extract_slide_title(slide) or f"Slide {slide_idx + 1}"
        slide_header = create_heading(title, level=1)
        slide_header.attributes = {"id": f"slide-{slide_idx + 1}"}
        slide_container.add_child(slide_header)
        
        # Add slide layout info as metadata callout
        layout_info = self._get_layout_info(slide)
        if layout_info:
            callout = self._create_obsidian_callout("info", "Slide Layout", layout_info)
            slide_container.add_child(callout)
        
        # Process slide content with categorization
        content_by_type = self._categorize_slide_content(slide)
        
        # Add main content
        for content_type, shapes in content_by_type.items():
            if shapes:
                section = self._process_content_section(content_type, shapes, slide_idx)
                if section:
                    slide_container.add_child(section)
        
        # Add speaker notes if present
        if slide.has_notes_slide:
            notes = self._extract_speaker_notes(slide)
            if notes:
                notes_section = self._create_collapsible_section("Speaker Notes", notes)
                slide_container.add_child(notes_section)
        
        # Add slide separator
        if slide_idx < len(prs.slides) - 1:
            separator = create_paragraph("---")
            slide_container.add_child(separator)
            
        return slide_container

    def _create_slide_navigation(self, current_idx: int, total_slides: int) -> Optional[ASTNode]:
        """Create navigation links between slides."""
        nav_parts = []
        
        # Previous slide link
        if current_idx > 0:
            prev_title = self.slide_titles.get(current_idx - 1, f"Slide {current_idx}")
            nav_parts.append(f"[[#slide-{current_idx}|← Previous]]")
        
        # Slide counter
        nav_parts.append(f"**{current_idx + 1} / {total_slides}**")
        
        # Next slide link
        if current_idx < total_slides - 1:
            next_title = self.slide_titles.get(current_idx + 1, f"Slide {current_idx + 2}")
            nav_parts.append(f"[[#slide-{current_idx + 2}|Next →]]")
        
        nav_text = " | ".join(nav_parts)
        nav_node = create_paragraph(nav_text)
        nav_node.style = StyleInfo()
        nav_node.style.text_align = "center"
        
        return nav_node

    def _get_layout_info(self, slide) -> Optional[str]:
        """Get information about the slide layout."""
        layout_name = slide.slide_layout.name if slide.slide_layout else "Custom"
        return f"Layout: {layout_name}"

    def _categorize_slide_content(self, slide) -> Dict[str, List]:
        """Categorize slide shapes by content type."""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            # Fallback if enum not available
            MSO_SHAPE_TYPE = type('MSO_SHAPE_TYPE', (), {
                'PICTURE': 13,
                'TABLE': 19,
                'GROUP': 6,
                'PLACEHOLDER': 14,
                'AUTO_SHAPE': 1
            })
        
        content = {
            "text": [],
            "images": [],
            "tables": [],
            "charts": [],
            "diagrams": [],
            "media": []
        }
        
        for shape in slide.shapes:
            shape_type = getattr(shape, 'shape_type', None)
            
            # Debug output for understanding shape types
            if hasattr(shape, 'image'):
                print(f"    Found image shape: type={shape_type}, has_image={hasattr(shape, 'image')}")
            
            # Check for images in multiple ways
            if shape_type == MSO_SHAPE_TYPE.PICTURE or shape_type == 13:
                content["images"].append(shape)
            elif hasattr(shape, 'image') and shape.image:
                # Sometimes images are in other shape types
                content["images"].append(shape)
            elif hasattr(shape, 'fill') and hasattr(shape.fill, 'type') and shape.fill.type == 6:  # Picture fill
                # Shape with picture fill
                content["images"].append(shape)
            elif shape_type == MSO_SHAPE_TYPE.TABLE or shape_type == 19:
                content["tables"].append(shape)
            elif hasattr(shape, 'has_table') and shape.has_table:
                content["tables"].append(shape)
            elif hasattr(shape, 'chart') and shape.chart:
                content["charts"].append(shape)
            elif hasattr(shape, 'has_chart') and shape.has_chart:
                content["charts"].append(shape)
            elif shape_type == MSO_SHAPE_TYPE.GROUP or shape_type == 6:
                # Group might contain images
                content["diagrams"].append(shape)
                # Also check for images within groups
                self._extract_images_from_group(shape, content["images"])
            elif hasattr(shape, 'media_type'):
                content["media"].append(shape)
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                if hasattr(shape.text_frame, 'text') and shape.text_frame.text.strip():
                    content["text"].append(shape)
                    
        # Report what we found
        if content["images"]:
            print(f"  Found {len(content['images'])} image(s) on slide")
                    
        return content
    
    def _extract_images_from_group(self, group_shape, images_list):
        """Extract images from grouped shapes."""
        try:
            if hasattr(group_shape, 'shapes'):
                for shape in group_shape.shapes:
                    if hasattr(shape, 'image') and shape.image:
                        images_list.append(shape)
                    elif hasattr(shape, 'shape_type') and shape.shape_type == 13:  # PICTURE
                        images_list.append(shape)
                    elif hasattr(shape, 'shapes'):  # Nested group
                        self._extract_images_from_group(shape, images_list)
        except Exception as e:
            print(f"    Warning: Could not extract images from group: {e}")

    def _process_content_section(self, content_type: str, shapes: List, slide_idx: int) -> Optional[ASTNode]:
        """Process a section of content based on its type."""
        if not shapes:
            return None
            
        section = ASTNode(NodeType.SECTION)
        
        if content_type == "text":
            for shape in shapes:
                text_node = self._process_text_shape_enhanced(shape)
                if text_node:
                    section.add_child(text_node)
                    
        elif content_type == "images":
            for shape in shapes:
                image_node = self._process_image_enhanced(shape, slide_idx)
                if image_node:
                    section.add_child(image_node)
                    
        elif content_type == "tables":
            for shape in shapes:
                table_node = self._process_table_enhanced(shape)
                if table_node:
                    section.add_child(table_node)
                    
        elif content_type == "charts":
            for shape in shapes:
                chart_node = self._process_chart_enhanced(shape, slide_idx)
                if chart_node:
                    section.add_child(chart_node)
                    
        elif content_type == "diagrams":
            for shape in shapes:
                diagram_node = self._process_diagram(shape, slide_idx)
                if diagram_node:
                    section.add_child(diagram_node)
                    
        elif content_type == "media":
            for shape in shapes:
                media_node = self._process_media(shape)
                if media_node:
                    section.add_child(media_node)
                    
        return section if section.children else None

    def _process_text_shape_enhanced(self, shape) -> Optional[ASTNode]:
        """Process text shape with enhanced formatting preservation."""
        if not shape.text_frame:
            return None
            
        text_frame = shape.text_frame
        
        # Check if it's a title or subtitle (safely)
        try:
            if shape.is_placeholder and hasattr(shape, 'placeholder_format'):
                ph_type = shape.placeholder_format.type
                if ph_type == 1:  # Title
                    return create_heading(text_frame.text, level=2)
                elif ph_type == 2:  # Subtitle
                    return create_heading(text_frame.text, level=3)
        except (ValueError, AttributeError):
            # Not a placeholder or error accessing placeholder_format
            pass
        
        # Process as formatted text
        formatted_content = []
        
        for paragraph in text_frame.paragraphs:
            para_parts = []
            
            # Check for bullet points
            if paragraph.level is not None:
                indent = "  " * paragraph.level
                para_parts.append(f"{indent}- ")
            
            for run in paragraph.runs:
                if run.text:
                    text = run.text
                    
                    # Apply formatting
                    if run.font.bold:
                        text = f"**{text}**"
                    if run.font.italic:
                        text = f"*{text}*"
                    if run.font.underline:
                        text = f"<u>{text}</u>"
                        
                    # Apply color if not default
                    if run.font.color and run.font.color.rgb:
                        rgb = run.font.color.rgb
                        # RGBColor is a tuple-like object (r, g, b)
                        if hasattr(rgb, '__iter__'):
                            try:
                                r, g, b = rgb
                                hex_color = f"#{r:02x}{g:02x}{b:02x}"
                                text = f'<span style="color: {hex_color}">{text}</span>'
                            except:
                                pass  # Skip color if format is unexpected
                        
                    para_parts.append(text)
            
            if para_parts:
                formatted_content.append(''.join(para_parts))
        
        if formatted_content:
            return create_paragraph('\n'.join(formatted_content))
        
        return None

    def _process_image_enhanced(self, shape, slide_idx: int) -> Optional[ASTNode]:
        """Extract and process image with Obsidian embedding."""
        try:
            # Check if this is actually an image shape
            if not hasattr(shape, 'image') or not shape.image:
                # Try to get the image part for picture shapes
                if hasattr(shape, 'click_action'):
                    # This might be a linked image
                    pass
                return create_paragraph(f"[Image placeholder - no image data found]")
            
            # Get the image part
            image_part = shape.image
            
            # Extract image blob
            image_stream = image_part.blob
            
            # Determine image format - be careful with accessing ext property
            ext = 'png'  # default
            try:
                # Try to get extension from the image part
                if hasattr(image_part, '_ext'):
                    ext = image_part._ext
                elif hasattr(image_part, 'content_type'):
                    content_type = image_part.content_type
                    if 'jpeg' in content_type or 'jpg' in content_type:
                        ext = 'jpg'
                    elif 'png' in content_type:
                        ext = 'png'
                    elif 'gif' in content_type:
                        ext = 'gif'
                    elif 'bmp' in content_type:
                        ext = 'bmp'
                    elif 'tiff' in content_type:
                        ext = 'tiff'
                    elif 'emf' in content_type or 'x-emf' in content_type:
                        ext = 'emf'
                    elif 'wmf' in content_type or 'x-wmf' in content_type:
                        ext = 'wmf'
                        
                # Try to detect format from the blob content
                if ext == 'png' and image_stream:
                    # Check magic bytes
                    if image_stream[:4] == b'\x89PNG':
                        ext = 'png'
                    elif image_stream[:2] in [b'\xff\xd8', b'JFIF']:
                        ext = 'jpg'
                    elif image_stream[:2] == b'BM':
                        ext = 'bmp'
                    elif image_stream[:3] == b'GIF':
                        ext = 'gif'
                    elif len(image_stream) > 40 and image_stream[40:44] == b' EMF':
                        ext = 'emf'
                    elif len(image_stream) > 0 and image_stream[:4] == b'\xd7\xcd\xc6\x9a':
                        ext = 'wmf'
            except Exception as e:
                print(f"    Warning: Could not determine image format: {e}")
                ext = 'png'  # fallback
            
            # Create image directory if needed
            if not self.image_dir:
                return create_paragraph(f"[Image: No output directory specified]")
                
            self.image_dir.mkdir(parents=True, exist_ok=True)
            
            # Save image
            self.image_counter += 1
            image_name = f"slide_{slide_idx + 1:02d}_image_{self.image_counter:03d}.{ext}"
            image_path = self.image_dir / image_name
            
            # Write the image data
            with open(image_path, 'wb') as f:
                f.write(image_stream)
            
            # Convert EMF/WMF to PNG if needed (for better Obsidian support)
            if ext in ['emf', 'wmf']:
                converted = self._convert_vector_to_png(image_path, image_stream)
                if converted:
                    # Update to use PNG version
                    old_path = image_path
                    image_name = image_name.replace(f'.{ext}', '.png')
                    image_path = self.image_dir / image_name
                    ext = 'png'
                    print(f"  ✅ Converted {old_path.name} to PNG format")
                    
            print(f"  ✅ Extracted image: {image_name} ({len(image_stream):,} bytes)")
            
            # Create Obsidian image embed
            alt_text = f"Image {self.image_counter} from slide {slide_idx + 1}"
            
            # Try to get better alt text from shape properties
            if hasattr(shape, 'name') and shape.name and not shape.name.startswith('Picture'):
                alt_text = shape.name
            elif hasattr(shape, 'alt_text') and shape.alt_text:
                alt_text = shape.alt_text
            
            # Get image dimensions for metadata
            width_inches = None
            height_inches = None
            if hasattr(shape, 'width') and hasattr(shape, 'height'):
                try:
                    from pptx.util import Inches
                    width_inches = shape.width / 914400  # EMUs to inches
                    height_inches = shape.height / 914400
                except:
                    pass
            
            image_node = ASTNode(NodeType.IMAGE)
            
            # Use Obsidian's internal link format for images
            # Use relative path from the markdown file location
            relative_path = f"{self.image_dir.name}/{image_name}"
            image_node.content = f"![[{relative_path}|{alt_text}]]"
            
            # Add metadata comment if we have dimensions
            container = ASTNode(NodeType.SECTION)
            container.add_child(image_node)
            
            if width_inches and height_inches:
                metadata = create_paragraph(f"<!-- Image dimensions: {width_inches:.1f}\" x {height_inches:.1f}\" -->")
                container.add_child(metadata)
            
            # Add caption if available
            if hasattr(shape, 'alt_text') and shape.alt_text:
                caption = create_paragraph(f"*{shape.alt_text}*")
                caption.style = StyleInfo()
                caption.style.text_align = "center"
                container.add_child(caption)
                return container
            
            return container if container.children else image_node
                
        except Exception as e:
            import traceback
            print(f"  ⚠️  Failed to extract image from slide {slide_idx + 1}: {e}")
            traceback.print_exc()
            # Fallback: create placeholder with error info
            return create_paragraph(f"[Image: Unable to extract - {type(e).__name__}: {str(e)}]")
    
    def _convert_vector_to_png(self, image_path: Path, image_data: bytes) -> bool:
        """Convert EMF/WMF to PNG format for better compatibility."""
        try:
            # Try using Pillow with additional support
            from PIL import Image
            import io
            
            # Try to open and convert
            try:
                img = Image.open(io.BytesIO(image_data))
                png_path = image_path.with_suffix('.png')
                img.save(png_path, 'PNG')
                # Remove original EMF/WMF file
                image_path.unlink()
                return True
            except:
                pass
                
            # Try using wand/ImageMagick if available
            try:
                from wand.image import Image as WandImage
                with WandImage(blob=image_data) as img:
                    img.format = 'png'
                    png_path = image_path.with_suffix('.png')
                    img.save(filename=str(png_path))
                    image_path.unlink()
                    return True
            except ImportError:
                pass
                
            # If we can't convert, keep the original
            return False
            
        except Exception as e:
            print(f"    Warning: Could not convert vector image: {e}")
            return False

    def _process_table_enhanced(self, shape) -> Optional[ASTNode]:
        """Process table with Obsidian-compatible formatting."""
        if not hasattr(shape, 'table'):
            return None
            
        table = shape.table
        table_node = ASTNode(NodeType.TABLE)
        
        for row_idx, row in enumerate(table.rows):
            row_node = ASTNode(NodeType.TABLE_ROW)
            
            for cell in row.cells:
                # Extract cell text with formatting
                cell_text = self._extract_cell_text_formatted(cell)
                
                # First row is header
                cell_type = NodeType.TABLE_HEADER if row_idx == 0 else NodeType.TABLE_CELL
                cell_node = ASTNode(cell_type, content=cell_text)
                
                # Add cell styling if present (safely)
                try:
                    if hasattr(cell, 'fill') and cell.fill and hasattr(cell.fill, 'fore_color') and cell.fill.fore_color:
                        if hasattr(cell.fill.fore_color, 'rgb') and cell.fill.fore_color.rgb:
                            rgb = cell.fill.fore_color.rgb
                            # RGBColor is a tuple-like object 
                            if hasattr(rgb, '__iter__'):
                                r, g, b = rgb
                                hex_color = f"#{r:02x}{g:02x}{b:02x}"
                                cell_node.style = StyleInfo()
                                cell_node.style.background_color = hex_color
                except (AttributeError, TypeError):
                    pass  # Skip cell styling if attributes don't exist
                
                row_node.add_child(cell_node)
            
            table_node.add_child(row_node)
        
        return table_node

    def _extract_cell_text_formatted(self, cell) -> str:
        """Extract formatted text from a table cell."""
        if not cell.text_frame:
            return cell.text if hasattr(cell, 'text') else ''
            
        parts = []
        for paragraph in cell.text_frame.paragraphs:
            para_text = []
            for run in paragraph.runs:
                if run.text:
                    text = run.text
                    if run.font.bold:
                        text = f"**{text}**"
                    if run.font.italic:
                        text = f"*{text}*"
                    para_text.append(text)
            if para_text:
                parts.append(''.join(para_text))
        
        return ' '.join(parts)

    def _process_chart_enhanced(self, shape, slide_idx: int) -> Optional[ASTNode]:
        """Process chart with data extraction and visualization description."""
        if not hasattr(shape, 'chart'):
            return None
            
        chart = shape.chart
        container = ASTNode(NodeType.SECTION)
        
        # Extract chart title
        title = "Chart"
        if chart.has_title and chart.chart_title:
            title = chart.chart_title.text_frame.text
        
        # Create descriptive callout
        chart_type = str(chart.chart_type).replace('_', ' ').title()
        description = f"**{title}**\nType: {chart_type}"
        
        # Try to extract data series information
        if hasattr(chart, 'series'):
            series_info = []
            for series in chart.series:
                if hasattr(series, 'name'):
                    series_info.append(f"- {series.name}")
            if series_info:
                description += "\n\nData Series:\n" + '\n'.join(series_info)
        
        callout = self._create_obsidian_callout("chart", "Chart Data", description)
        container.add_child(callout)
        
        # Try to save chart as image if possible
        try:
            # This would require additional libraries for chart rendering
            # For now, we'll just add a placeholder
            placeholder = create_paragraph(f"[Chart visualization: {title}]")
            container.add_child(placeholder)
        except:
            pass
        
        return container

    def _process_diagram(self, shape, slide_idx: int) -> Optional[ASTNode]:
        """Process grouped shapes as a diagram."""
        container = ASTNode(NodeType.SECTION)
        
        # Add diagram indicator
        diagram_note = self._create_obsidian_callout(
            "abstract", 
            "Diagram", 
            "Complex diagram or grouped shapes detected"
        )
        container.add_child(diagram_note)
        
        # Try to extract text from grouped shapes
        text_content = self._extract_text_from_group(shape)
        if text_content:
            for text in text_content:
                container.add_child(create_paragraph(text))
        
        return container

    def _extract_text_from_group(self, group_shape) -> List[str]:
        """Recursively extract text from grouped shapes."""
        texts = []
        
        if hasattr(group_shape, 'shapes'):
            for shape in group_shape.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    if shape.text_frame.text.strip():
                        texts.append(shape.text_frame.text.strip())
                elif hasattr(shape, 'shapes'):  # Nested group
                    texts.extend(self._extract_text_from_group(shape))
        
        return texts

    def _process_media(self, shape) -> Optional[ASTNode]:
        """Process embedded media (video/audio)."""
        media_type = "Media"
        if hasattr(shape, 'media_type'):
            media_type = str(shape.media_type).replace('_', ' ').title()
        
        return self._create_obsidian_callout(
            "example",
            f"{media_type} Content",
            f"Embedded {media_type.lower()} file detected. Manual extraction may be required."
        )

    def _extract_speaker_notes(self, slide) -> Optional[str]:
        """Extract speaker notes from slide."""
        if not slide.has_notes_slide:
            return None
            
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ''
        
        return notes_text.strip() if notes_text else None

    def _create_collapsible_section(self, title: str, content: str) -> ASTNode:
        """Create a collapsible section using Obsidian syntax."""
        section = ASTNode(NodeType.SECTION)
        
        # Use HTML details tag for collapsible content
        details_content = f"<details>\n<summary>{title}</summary>\n\n{content}\n\n</details>"
        section.add_child(create_paragraph(details_content))
        
        return section

    def _create_obsidian_callout(self, callout_type: str, title: str, content: str) -> ASTNode:
        """Create an Obsidian-style callout block."""
        # Obsidian callout syntax: > [!type] Title
        # > Content
        
        callout_lines = [
            f"> [!{callout_type}] {title}",
            *[f"> {line}" for line in content.split('\n')]
        ]
        
        return create_paragraph('\n'.join(callout_lines))

    def _create_presentation_summary(self, prs) -> Optional[ASTNode]:
        """Create a summary section for the presentation."""
        summary = ASTNode(NodeType.SECTION)
        
        # Add summary heading
        summary.add_child(create_heading("Presentation Summary", level=1))
        
        # Add statistics
        stats = [
            f"- Total Slides: {len(prs.slides)}",
            f"- Extracted Images: {self.image_counter}",
            f"- Slide Titles: {len(self.slide_titles)}"
        ]
        
        summary.add_child(create_paragraph('\n'.join(stats)))
        
        # Add table of contents
        if self.slide_titles:
            toc_heading = create_heading("Table of Contents", level=2)
            summary.add_child(toc_heading)
            
            toc_items = []
            for idx, title in self.slide_titles.items():
                toc_items.append(f"- [[#slide-{idx + 1}|Slide {idx + 1}: {title}]]")
            
            summary.add_child(create_paragraph('\n'.join(toc_items)))
        
        return summary

    def ast2md_obsidian(self, ast_root: ASTNode, output_path: str) -> None:
        """Convert AST to Obsidian-optimized Markdown format."""
        content_parts = []
        
        # Add YAML frontmatter with metadata
        if ast_root.metadata:
            frontmatter = self._create_yaml_frontmatter(ast_root.metadata)
            content_parts.append(frontmatter)
            content_parts.append("")  # Empty line after frontmatter
        
        # Convert AST to Markdown with Obsidian features
        markdown_content = self._ast_to_markdown_obsidian(ast_root)
        content_parts.append(markdown_content)
        
        # Write to file
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(content_parts))

    def _create_yaml_frontmatter(self, metadata: DocumentMetadata) -> str:
        """Create YAML frontmatter for Obsidian."""
        yaml_lines = ["---"]
        
        # Standard metadata
        if metadata.title:
            yaml_lines.append(f"title: \"{metadata.title}\"")
        if metadata.author:
            yaml_lines.append(f"author: \"{metadata.author}\"")
        if metadata.subject:
            yaml_lines.append(f"subject: \"{metadata.subject}\"")
        if metadata.keywords:
            yaml_lines.append(f"keywords: \"{metadata.keywords}\"")
            
        # Dates
        if metadata.created:
            yaml_lines.append(f"created: {metadata.created.isoformat() if hasattr(metadata.created, 'isoformat') else metadata.created}")
        if metadata.modified:
            yaml_lines.append(f"modified: {metadata.modified.isoformat() if hasattr(metadata.modified, 'isoformat') else metadata.modified}")
            
        # Custom Obsidian metadata
        if metadata.custom:
            if "tags" in metadata.custom:
                tags = metadata.custom["tags"]
                yaml_lines.append(f"tags: [{', '.join(tags)}]")
            if "aliases" in metadata.custom:
                aliases = metadata.custom["aliases"]
                yaml_lines.append(f"aliases: [{', '.join(f'\"{a}\"' for a in aliases)}]")
            if "total_slides" in metadata.custom:
                yaml_lines.append(f"total_slides: {metadata.custom['total_slides']}")
            if "conversion_date" in metadata.custom:
                yaml_lines.append(f"conversion_date: {metadata.custom['conversion_date']}")
        
        yaml_lines.append("---")
        
        return '\n'.join(yaml_lines)

    def _ast_to_markdown_obsidian(self, ast_root: ASTNode) -> str:
        """Convert AST to Markdown with Obsidian-specific features."""
        lines = []
        
        def process_node(node: ASTNode, indent_level: int = 0):
            indent = "  " * indent_level
            
            if node.type == NodeType.HEADING:
                level = node.attributes.get('level', 1) if node.attributes else 1
                heading_id = node.attributes.get('id', '') if node.attributes else ''
                heading_text = node.content or ''
                
                # Add anchor if ID is present
                if heading_id:
                    lines.append(f"{'#' * level} {heading_text} ^{heading_id}")
                else:
                    lines.append(f"{'#' * level} {heading_text}")
                    
            elif node.type == NodeType.PARAGRAPH:
                if node.content:
                    lines.append(f"{indent}{node.content}")
                    
            elif node.type == NodeType.TABLE:
                table_lines = self._render_table_obsidian(node)
                lines.extend(table_lines)
                
            elif node.type == NodeType.IMAGE:
                if node.content:
                    lines.append(node.content)  # Already formatted for Obsidian
                elif node.attributes:
                    src = node.attributes.get('src', '')
                    alt = node.attributes.get('alt', '')
                    lines.append(f"![[{src}|{alt}]]")
                    
            elif node.type in [NodeType.LIST_ORDERED, NodeType.LIST_UNORDERED]:
                for idx, child in enumerate(node.children, 1):
                    if child.type == NodeType.LIST_ITEM:
                        if node.type == NodeType.LIST_ORDERED:
                            lines.append(f"{indent}{idx}. {child.content or ''}")
                        else:
                            lines.append(f"{indent}- {child.content or ''}")
                        
            elif node.type == NodeType.CODE_BLOCK:
                language = node.attributes.get('language', '') if node.attributes else ''
                lines.append(f"```{language}")
                if node.content:
                    lines.append(node.content)
                lines.append("```")
                
            elif node.type == NodeType.LINE_BREAK:
                lines.append("")
                
            # Process children
            for child in node.children:
                process_node(child, indent_level)
        
        # Process all nodes
        for child in ast_root.children:
            process_node(child)
        
        return '\n'.join(lines)

    def _render_table_obsidian(self, table_node: ASTNode) -> List[str]:
        """Render table in Obsidian-compatible Markdown format."""
        lines = []
        
        # Process rows
        for row_idx, row in enumerate(table_node.children):
            if row.type == NodeType.TABLE_ROW:
                cells = []
                for cell in row.children:
                    content = cell.content or ''
                    cells.append(content)
                
                # Create row
                lines.append(f"| {' | '.join(cells)} |")
                
                # Add header separator after first row
                if row_idx == 0:
                    separators = ['-' * max(len(cell), 3) for cell in cells]
                    lines.append(f"| {' | '.join(separators)} |")
        
        return lines

    def convert(self, input_path: str, output_path: str) -> None:
        """Convert PPTX to Obsidian-optimized Markdown."""
        
        try:
            # Parse PPTX to enhanced AST
            ast_root = self.parse_pptx2ast(input_path)
            
            # Convert to Obsidian Markdown
            self.ast2md_obsidian(ast_root, output_path)
            
            # Create companion note with conversion info
            self._create_companion_note(input_path, output_path)
            
            print(f"Successfully converted '{input_path}' to Obsidian Markdown '{output_path}'")
            
            if self.image_dir and self.image_counter > 0:
                print(f"Extracted {self.image_counter} images to '{self.image_dir}'")
                
        except Exception as e:
            raise ConversionError(
                f"PPTX to Obsidian Markdown conversion failed: {e}",
                source_format="pptx",
                target_format="md",
                suggestions=[
                    "Ensure python-pptx is installed: pip install python-pptx",
                    "Check that the PPTX file is valid and not corrupted",
                    "Verify write permissions for output directory"
                ]
            )

    def _create_companion_note(self, input_path: str, output_path: str):
        """Create a companion note with conversion metadata."""
        companion_path = Path(output_path).parent / f"{Path(output_path).stem}_info.md"
        
        info_content = f"""---
title: "Conversion Info: {Path(input_path).stem}"
tags: [conversion-info, pptx-to-md]
created: {datetime.now().isoformat()}
---

# Conversion Information

## Source File
- **Path**: `{input_path}`
- **Format**: PowerPoint Presentation (PPTX)

## Output File
- **Path**: `{output_path}`
- **Format**: Obsidian-optimized Markdown

## Conversion Details
- **Converter**: Pptx2MdObsidianConverter
- **Date**: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
- **Images Extracted**: {self.image_counter}
- **Image Directory**: `{self.image_dir.name if self.image_dir else 'N/A'}`

## Features Applied
- ✅ YAML frontmatter with metadata
- ✅ Slide navigation links
- ✅ Image extraction and embedding
- ✅ Table formatting preservation
- ✅ Speaker notes as collapsible sections
- ✅ Obsidian callouts for special content
- ✅ Internal links and cross-references
- ✅ Formatted text preservation (bold, italic, colors)
- ✅ Table of contents generation

## Usage Tips
1. Place extracted images folder in your Obsidian vault
2. Use graph view to see slide connections
3. Enable slide view plugin for presentation mode
4. Use tags for organizing multiple presentations
"""
        
        with open(companion_path, 'w', encoding='utf-8') as f:
            f.write(info_content)