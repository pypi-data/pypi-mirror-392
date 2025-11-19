#!/usr/bin/env python3
"""
Visual HTML to PPTX converter that preserves exact layout using browser rendering.
This converter uses browser automation to capture the visual representation of HTML slides
and embeds them as high-quality images in PowerPoint slides.
"""

import os
import re
import sys
import asyncio
import tempfile
import subprocess
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from dataclasses import dataclass

try:
    from playwright.async_api import async_playwright, Browser, Page
    PLAYWRIGHT_AVAILABLE = True
except ImportError:
    PLAYWRIGHT_AVAILABLE = False

try:
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    from selenium.webdriver.common.by import By
    from selenium.webdriver.support.ui import WebDriverWait
    from selenium.webdriver.support import expected_conditions as EC
    SELENIUM_AVAILABLE = True
except ImportError:
    SELENIUM_AVAILABLE = False

from bs4 import BeautifulSoup, Tag
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from ..core.exceptions import ConversionError


@dataclass
class SlideCapture:
    """Data class for captured slide information."""
    index: int
    title: str
    image_path: str
    width: int
    height: int
    original_element: Optional[Tag] = None


class VisualHtml2PptxConverter:
    """
    Advanced HTML to PPTX converter that uses browser automation to capture
    exact visual representation of HTML slides.
    """
    
    def __init__(self, 
                 target_width: int = 1920, 
                 target_height: int = 1080,
                 quality: str = 'high',
                 browser_engine: str = 'auto'):
        """
        Initialize the visual converter.
        
        Args:
            target_width: Target slide width in pixels
            target_height: Target slide height in pixels  
            quality: Image quality ('low', 'medium', 'high', 'ultra')
            browser_engine: Browser to use ('playwright', 'selenium', 'auto')
        """
        self.target_width = target_width
        self.target_height = target_height
        self.quality = quality
        self.browser_engine = browser_engine
        
        # Quality settings
        self.quality_settings = {
            'low': {'scale': 1, 'format': 'jpeg', 'jpeg_quality': 70},
            'medium': {'scale': 1.5, 'format': 'png', 'jpeg_quality': 85},
            'high': {'scale': 2, 'format': 'png', 'jpeg_quality': 95},
            'ultra': {'scale': 3, 'format': 'png', 'jpeg_quality': 100}
        }
        
        self.current_settings = self.quality_settings.get(quality, self.quality_settings['high'])
        
        # Temporary directory for captures
        self.temp_dir = None
        
    async def convert(self, input_path: str, output_path: str) -> None:
        """
        Convert HTML presentation to PPTX with visual fidelity.
        
        Args:
            input_path: Path to HTML file
            output_path: Path for output PPTX file
        """
        
        if not os.path.exists(input_path):
            raise ConversionError(f"Input file not found: {input_path}")
        
        print(f"Starting visual HTML to PPTX conversion...")
        print(f"Quality: {self.quality} ({self.current_settings})")
        
        # Create temporary directory
        self.temp_dir = tempfile.mkdtemp(prefix="html2pptx_")
        
        try:
            # Choose browser engine
            if self.browser_engine == 'auto':
                if PLAYWRIGHT_AVAILABLE:
                    captures = await self._capture_with_playwright(input_path)
                elif SELENIUM_AVAILABLE:
                    captures = await self._capture_with_selenium(input_path)
                else:
                    captures = await self._capture_with_fallback(input_path)
            elif self.browser_engine == 'playwright':
                captures = await self._capture_with_playwright(input_path)
            elif self.browser_engine == 'selenium':
                captures = await self._capture_with_selenium(input_path)
            else:
                raise ConversionError(f"Unknown browser engine: {self.browser_engine}")
            
            if not captures:
                raise ConversionError("No slides were captured from the HTML file")
            
            # Create PowerPoint presentation
            await self._create_pptx(captures, output_path)
            
            print(f"Successfully converted {len(captures)} slides to {output_path}")
            
        finally:
            # Cleanup temporary files
            self._cleanup_temp_files()
    
    async def _capture_with_playwright(self, input_path: str) -> List[SlideCapture]:
        """Capture slides using Playwright browser automation."""
        
        if not PLAYWRIGHT_AVAILABLE:
            raise ConversionError("Playwright not available. Install with: pip install playwright")
        
        print("Using Playwright for high-quality capture...")
        captures = []
        
        async with async_playwright() as p:
            # Launch browser with high DPI support
            browser = await p.chromium.launch(
                headless=True,
                args=[
                    '--force-device-scale-factor=' + str(self.current_settings['scale']),
                    '--high-dpi-support=1',
                    '--force-color-profile=srgb',
                    '--disable-dev-shm-usage',
                    '--no-sandbox'
                ]
            )
            
            try:
                page = await browser.new_page(
                    viewport={
                        'width': self.target_width,
                        'height': self.target_height
                    },
                    device_scale_factor=self.current_settings['scale']
                )
                
                # Load HTML file
                file_url = f"file://{os.path.abspath(input_path)}"
                await page.goto(file_url, wait_until='networkidle')
                
                # Wait for any dynamic content
                await page.wait_for_timeout(2000)
                
                # Find slide elements
                slides = await self._detect_slides_in_page(page)
                
                print(f"Detected {len(slides)} slides for capture")
                
                for i, slide_element in enumerate(slides):
                    print(f"Capturing slide {i + 1}/{len(slides)}...")
                    
                    # Show only this slide
                    await self._show_slide(page, i)
                    await page.wait_for_timeout(500)
                    
                    # Capture slide
                    capture_path = os.path.join(self.temp_dir, f"slide_{i:03d}.{self.current_settings['format']}")
                    
                    await page.screenshot(
                        path=capture_path,
                        full_page=False,
                        type=self.current_settings['format'],
                        quality=self.current_settings.get('jpeg_quality', 95) if self.current_settings['format'] == 'jpeg' else None
                    )
                    
                    # Get slide title
                    title = await self._extract_slide_title(page, slide_element)
                    
                    # Create capture record
                    captures.append(SlideCapture(
                        index=i,
                        title=title,
                        image_path=capture_path,
                        width=self.target_width,
                        height=self.target_height
                    ))
                
            finally:
                await browser.close()
        
        return captures
    
    async def _capture_with_selenium(self, input_path: str) -> List[SlideCapture]:
        """Capture slides using Selenium WebDriver."""
        
        if not SELENIUM_AVAILABLE:
            raise ConversionError("Selenium not available. Install with: pip install selenium")
        
        print("Using Selenium for slide capture...")
        captures = []
        
        # Configure Chrome options
        chrome_options = Options()
        chrome_options.add_argument('--headless')
        chrome_options.add_argument('--no-sandbox')
        chrome_options.add_argument('--disable-dev-shm-usage')
        chrome_options.add_argument(f'--window-size={self.target_width},{self.target_height}')
        chrome_options.add_argument(f'--force-device-scale-factor={self.current_settings["scale"]}')
        
        driver = webdriver.Chrome(options=chrome_options)
        
        try:
            # Load HTML file
            file_url = f"file://{os.path.abspath(input_path)}"
            driver.get(file_url)
            
            # Wait for page load
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.TAG_NAME, "body"))
            )
            
            # Find slides
            slide_elements = driver.find_elements(By.CSS_SELECTOR, ".slide")
            
            if not slide_elements:
                # Try alternative selectors
                slide_elements = driver.find_elements(By.CSS_SELECTOR, "[class*='slide']")
            
            print(f"Detected {len(slide_elements)} slides for capture")
            
            for i, slide_element in enumerate(slide_elements):
                print(f"Capturing slide {i + 1}/{len(slide_elements)}...")
                
                # Show slide (execute JavaScript)
                driver.execute_script(f"showSlide({i});")
                
                # Wait for animation
                driver.implicitly_wait(1)
                
                # Capture screenshot
                capture_path = os.path.join(self.temp_dir, f"slide_{i:03d}.png")
                driver.save_screenshot(capture_path)
                
                # Get slide title
                title_elements = slide_element.find_elements(By.TAG_NAME, "h1")
                if not title_elements:
                    title_elements = slide_element.find_elements(By.TAG_NAME, "h2")
                
                title = title_elements[0].text if title_elements else f"Slide {i + 1}"
                
                # Create capture record
                captures.append(SlideCapture(
                    index=i,
                    title=title,
                    image_path=capture_path,
                    width=self.target_width,
                    height=self.target_height
                ))
                
        finally:
            driver.quit()
        
        return captures
    
    async def _capture_with_fallback(self, input_path: str) -> List[SlideCapture]:
        """Fallback capture method using HTML parsing."""
        
        print("Using fallback method (limited visual fidelity)...")
        
        # Parse HTML
        with open(input_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        soup = BeautifulSoup(content, 'html.parser')
        slides = soup.find_all(class_='slide')
        
        if not slides:
            slides = soup.find_all(lambda tag: tag.get('class') and 
                                 any('slide' in cls.lower() for cls in tag.get('class', [])))
        
        captures = []
        
        for i, slide in enumerate(slides):
            print(f"Processing slide {i + 1}/{len(slides)} (text-based)...")
            
            # Create a basic image representation
            capture_path = os.path.join(self.temp_dir, f"slide_{i:03d}.png")
            self._create_text_slide_image(slide, capture_path, i)
            
            # Get title
            title_tag = slide.find(['h1', 'h2', 'h3'])
            title = title_tag.get_text(strip=True) if title_tag else f"Slide {i + 1}"
            
            captures.append(SlideCapture(
                index=i,
                title=title,
                image_path=capture_path,
                width=self.target_width,
                height=self.target_height,
                original_element=slide
            ))
        
        return captures
    
    async def _detect_slides_in_page(self, page: Page) -> List[Any]:
        """Detect slide elements in the page."""
        
        # Try different slide selectors
        slide_selectors = [
            '.slide',
            'section.slide',
            '[data-slide]',
            '.reveal .slides section',
            '.impress .step'
        ]
        
        for selector in slide_selectors:
            slides = await page.query_selector_all(selector)
            if slides:
                print(f"Found slides using selector: {selector}")
                return slides
        
        return []
    
    async def _show_slide(self, page: Page, slide_index: int) -> None:
        """Show specific slide in the presentation."""
        
        # Try common slide navigation functions
        try:
            await page.evaluate(f"showSlide({slide_index})")
        except:
            try:
                await page.evaluate(f"goToSlide({slide_index})")
            except:
                try:
                    # Hide all slides, show target slide
                    await page.evaluate(f"""
                        document.querySelectorAll('.slide').forEach((slide, index) => {{
                            slide.style.display = index === {slide_index} ? 'block' : 'none';
                        }});
                    """)
                except:
                    pass
    
    async def _extract_slide_title(self, page: Page, slide_element: Any) -> str:
        """Extract title from slide element."""
        
        try:
            # Try to get h1, h2, or h3 text
            title = await slide_element.query_selector('h1, h2, h3')
            if title:
                return await title.text_content()
        except:
            pass
        
        return "Untitled Slide"
    
    def _create_text_slide_image(self, slide_element: Tag, output_path: str, slide_index: int) -> None:
        """Create a basic image representation of slide content (fallback)."""
        
        # Create blank image
        img = Image.new('RGB', (self.target_width, self.target_height), 'white')
        draw = ImageDraw.Draw(img)
        
        try:
            # Try to load a font
            font_large = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 48)
            font_medium = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 32)
            font_small = ImageFont.truetype("/System/Library/Fonts/Arial.ttf", 24)
        except:
            # Fallback to default font
            font_large = ImageFont.load_default()
            font_medium = ImageFont.load_default()
            font_small = ImageFont.load_default()
        
        y_offset = 100
        
        # Draw title
        title_tag = slide_element.find(['h1', 'h2', 'h3'])
        if title_tag:
            title = title_tag.get_text(strip=True)
            draw.text((100, y_offset), title, fill='black', font=font_large)
            y_offset += 100
        
        # Draw content
        content_elements = slide_element.find_all(['p', 'li', 'div'])
        for element in content_elements[:10]:  # Limit to prevent overflow
            text = element.get_text(strip=True)
            if text and len(text) > 3:
                # Wrap text to fit
                words = text.split()
                lines = []
                current_line = []
                
                for word in words:
                    current_line.append(word)
                    line_text = ' '.join(current_line)
                    bbox = draw.textbbox((0, 0), line_text, font=font_small)
                    if bbox[2] > self.target_width - 200:  # Leave margin
                        if len(current_line) > 1:
                            current_line.pop()
                            lines.append(' '.join(current_line))
                            current_line = [word]
                        else:
                            lines.append(word)
                            current_line = []
                
                if current_line:
                    lines.append(' '.join(current_line))
                
                for line in lines[:3]:  # Max 3 lines per element
                    draw.text((120, y_offset), line, fill='black', font=font_small)
                    y_offset += 35
                
                y_offset += 20
                
                if y_offset > self.target_height - 100:
                    break
        
        # Add slide number
        draw.text((self.target_width - 100, self.target_height - 50), 
                 f"Slide {slide_index + 1}", fill='gray', font=font_small)
        
        img.save(output_path)
    
    async def _create_pptx(self, captures: List[SlideCapture], output_path: str) -> None:
        """Create PowerPoint presentation from captured slides."""
        
        print(f"Creating PowerPoint with {len(captures)} captured slides...")
        
        # Create presentation
        prs = Presentation()
        
        # Remove default slide
        if len(prs.slides) > 0:
            slide_part = prs.slides._sldIdLst[0]
            prs.part.drop_rel(slide_part.rId)
            del prs.slides._sldIdLst[0]
        
        # Add captured slides
        for capture in captures:
            print(f"Adding slide {capture.index + 1}: {capture.title}")
            
            # Use blank layout
            blank_layout = prs.slide_layouts[6]
            slide = prs.slides.add_slide(blank_layout)
            
            # Add image to fill slide
            slide.shapes.add_picture(
                capture.image_path,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            # Set slide title in notes
            if hasattr(slide, 'notes_slide'):
                slide.notes_slide.notes_text_frame.text = capture.title
        
        # Save presentation
        prs.save(output_path)
        print(f"Saved presentation to {output_path}")
    
    def _cleanup_temp_files(self) -> None:
        """Clean up temporary files."""
        
        if self.temp_dir and os.path.exists(self.temp_dir):
            import shutil
            try:
                shutil.rmtree(self.temp_dir)
                print("Cleaned up temporary files")
            except Exception as e:
                print(f"Warning: Could not clean up temporary files: {e}")


# Async wrapper for synchronous usage
def html2pptx_visual_convert(input_path: str, output_path: str, **kwargs) -> None:
    """Synchronous wrapper for visual HTML to PPTX conversion."""
    
    converter = VisualHtml2PptxConverter(**kwargs)
    
    # Run async conversion
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    
    try:
        loop.run_until_complete(converter.convert(input_path, output_path))
    finally:
        loop.close()


# Main converter class for compatibility
class Html2PptxVisualConverter:
    """Main converter class for visual HTML to PPTX conversion."""
    
    def __init__(self, **kwargs):
        self.kwargs = kwargs
    
    def convert(self, input_path: str, output_path: str) -> None:
        """Convert HTML to PPTX with visual fidelity."""
        html2pptx_visual_convert(input_path, output_path, **self.kwargs)