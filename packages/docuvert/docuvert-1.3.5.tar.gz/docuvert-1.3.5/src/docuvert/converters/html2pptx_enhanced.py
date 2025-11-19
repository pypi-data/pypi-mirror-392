"""
Enhanced HTML to PPTX converter with advanced layout analysis.

This is a specialized version for complex presentations with:
- CSS Grid and Flexbox layout analysis
- Precise positioning and sizing
- Advanced styling preservation
- Complex element relationships
"""

import sys
import os
import re
import json
import math
from typing import Optional, Dict, Any, List, Tuple, NamedTuple
from dataclasses import dataclass
from enum import Enum

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from .html2pptx import Html2PptxConverter, CSSStyleParser, AdvancedImageRenderer
from core.exceptions import ConversionError

try:
    from bs4 import BeautifulSoup, Tag
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_FILL_TYPE
    from pptx.shapes.shapetree import SlideShapes
    from pptx.shapes.autoshape import Shape
except ImportError as e:
    raise ConversionError(f"Required libraries missing: {e}")


class LayoutType(Enum):
    """Types of CSS layouts detected."""
    FLEX = "flex"
    GRID = "grid"
    ABSOLUTE = "absolute"
    RELATIVE = "relative"
    STATIC = "static"
    FLOAT = "float"


@dataclass
class ElementBounds:
    """Represents element positioning and dimensions."""
    x: float
    y: float
    width: float
    height: float
    z_index: int = 0
    
    def to_inches(self, scale: float = 1.0) -> Tuple[Inches, Inches, Inches, Inches]:
        """Convert to PowerPoint Inches with scaling."""
        px_to_inches = 1.0 / 96.0 * scale  # Assume 96 DPI with scaling
        return (
            Inches(self.x * px_to_inches),
            Inches(self.y * px_to_inches), 
            Inches(self.width * px_to_inches),
            Inches(self.height * px_to_inches)
        )


@dataclass
class ElementStyle:
    """Comprehensive element styling information."""
    # Text properties
    font_family: Optional[str] = None
    font_size: Optional[float] = None
    font_weight: Optional[str] = None
    font_style: Optional[str] = None
    color: Optional[Tuple[int, int, int]] = None
    text_align: Optional[str] = None
    line_height: Optional[float] = None
    
    # Background properties
    background_color: Optional[Tuple[int, int, int]] = None
    background_gradient: Optional[Dict[str, Any]] = None
    background_image: Optional[str] = None
    
    # Border properties
    border_width: Optional[float] = None
    border_color: Optional[Tuple[int, int, int]] = None
    border_style: Optional[str] = None
    border_radius: Optional[float] = None
    
    # Layout properties
    padding: Tuple[float, float, float, float] = (0, 0, 0, 0)  # top, right, bottom, left
    margin: Tuple[float, float, float, float] = (0, 0, 0, 0)
    
    # Effects
    opacity: float = 1.0
    box_shadow: Optional[str] = None
    transform: Optional[str] = None


class AdvancedCSSParser(CSSStyleParser):
    """Enhanced CSS parser with layout analysis capabilities."""
    
    def __init__(self):
        super().__init__()
        self.viewport_width = 1920
        self.viewport_height = 1080
    
    def parse_element_bounds(self, element: Tag, parent_bounds: Optional[ElementBounds] = None) -> ElementBounds:
        """Parse element positioning and dimensions from CSS."""
        style = element.get('style', '')
        computed_style = self._get_computed_styles(element, style)
        
        # Default bounds
        bounds = ElementBounds(0, 0, 100, 50)
        
        # Parse positioning
        position = computed_style.get('position', 'static')
        
        if position == 'absolute':
            # Absolute positioning
            left = self._parse_dimension(computed_style.get('left', '0'), self.viewport_width)
            top = self._parse_dimension(computed_style.get('top', '0'), self.viewport_height)
            bounds.x = left
            bounds.y = top
            
        elif position == 'relative' and parent_bounds:
            # Relative to parent
            left_offset = self._parse_dimension(computed_style.get('left', '0'), parent_bounds.width)
            top_offset = self._parse_dimension(computed_style.get('top', '0'), parent_bounds.height)
            bounds.x = parent_bounds.x + left_offset
            bounds.y = parent_bounds.y + top_offset
            
        elif parent_bounds:
            # Default positioning within parent
            bounds.x = parent_bounds.x
            bounds.y = parent_bounds.y
        
        # Parse dimensions
        width = self._parse_dimension(computed_style.get('width', 'auto'), self.viewport_width)
        height = self._parse_dimension(computed_style.get('height', 'auto'), self.viewport_height)
        
        if width > 0:
            bounds.width = width
        if height > 0:
            bounds.height = height
        
        # Parse z-index
        z_index = computed_style.get('z-index')
        if z_index and z_index.isdigit():
            bounds.z_index = int(z_index)
        
        return bounds
    
    def parse_comprehensive_styles(self, element: Tag) -> ElementStyle:
        """Parse comprehensive styling information from element."""
        style = element.get('style', '')
        computed_style = self._get_computed_styles(element, style)
        
        element_style = ElementStyle()
        
        # Font properties
        element_style.font_family = computed_style.get('font-family')
        if computed_style.get('font-size'):
            element_style.font_size = self.parse_font_size(computed_style['font-size'])
        element_style.font_weight = computed_style.get('font-weight')
        element_style.font_style = computed_style.get('font-style')
        element_style.text_align = computed_style.get('text-align')
        
        # Colors
        if computed_style.get('color'):
            element_style.color = self.parse_color(computed_style['color'])
        if computed_style.get('background-color'):
            element_style.background_color = self.parse_color(computed_style['background-color'])
        
        # Background gradient
        if computed_style.get('background'):
            gradient = self.parse_gradient(computed_style['background'])
            if gradient:
                element_style.background_gradient = gradient
        
        # Border properties
        if computed_style.get('border-width'):
            element_style.border_width = self._parse_dimension(computed_style['border-width'], 1)
        if computed_style.get('border-color'):
            element_style.border_color = self.parse_color(computed_style['border-color'])
        element_style.border_style = computed_style.get('border-style')
        if computed_style.get('border-radius'):
            element_style.border_radius = self._parse_dimension(computed_style['border-radius'], 1)
        
        # Spacing
        element_style.padding = self._parse_box_dimensions(computed_style, 'padding')
        element_style.margin = self._parse_box_dimensions(computed_style, 'margin')
        
        # Effects
        if computed_style.get('opacity'):
            try:
                element_style.opacity = float(computed_style['opacity'])
            except ValueError:
                pass
        
        element_style.box_shadow = computed_style.get('box-shadow')
        element_style.transform = computed_style.get('transform')
        
        return element_style
    
    def detect_layout_type(self, element: Tag) -> LayoutType:
        """Detect the layout type used by an element."""
        style = element.get('style', '')
        computed_style = self._get_computed_styles(element, style)
        
        display = computed_style.get('display', 'block')
        position = computed_style.get('position', 'static')
        
        if 'flex' in display:
            return LayoutType.FLEX
        elif 'grid' in display:
            return LayoutType.GRID
        elif position == 'absolute':
            return LayoutType.ABSOLUTE
        elif position == 'relative':
            return LayoutType.RELATIVE
        elif computed_style.get('float') not in [None, 'none']:
            return LayoutType.FLOAT
        else:
            return LayoutType.STATIC
    
    def analyze_flexbox_layout(self, container: Tag) -> Dict[str, Any]:
        """Analyze flexbox layout properties."""
        style = container.get('style', '')
        computed_style = self._get_computed_styles(container, style)
        
        return {
            'direction': computed_style.get('flex-direction', 'row'),
            'wrap': computed_style.get('flex-wrap', 'nowrap'),
            'justify_content': computed_style.get('justify-content', 'flex-start'),
            'align_items': computed_style.get('align-items', 'stretch'),
            'align_content': computed_style.get('align-content', 'stretch'),
            'gap': self._parse_dimension(computed_style.get('gap', '0'), 1)
        }
    
    def analyze_grid_layout(self, container: Tag) -> Dict[str, Any]:
        """Analyze CSS Grid layout properties."""
        style = container.get('style', '')
        computed_style = self._get_computed_styles(container, style)
        
        return {
            'template_columns': computed_style.get('grid-template-columns', 'none'),
            'template_rows': computed_style.get('grid-template-rows', 'none'),
            'gap': self._parse_dimension(computed_style.get('grid-gap', computed_style.get('gap', '0')), 1),
            'column_gap': self._parse_dimension(computed_style.get('grid-column-gap', '0'), 1),
            'row_gap': self._parse_dimension(computed_style.get('grid-row-gap', '0'), 1)
        }
    
    def _get_computed_styles(self, element: Tag, inline_style: str) -> Dict[str, str]:
        """Extract computed styles from inline styles and classes."""
        styles = {}
        
        # Parse inline styles
        for declaration in inline_style.split(';'):
            if ':' in declaration:
                property_name, value = declaration.split(':', 1)
                styles[property_name.strip()] = value.strip()
        
        # TODO: Parse CSS from stylesheets based on element classes/IDs
        # This would require a full CSS engine, which is complex
        
        return styles
    
    def _parse_dimension(self, value: str, reference_size: float = 1) -> float:
        """Parse CSS dimension value to pixels."""
        if not value or value == 'auto':
            return 0
        
        value = value.strip().lower()
        
        # Pixels
        if value.endswith('px'):
            try:
                return float(value[:-2])
            except ValueError:
                return 0
        
        # Percentages
        if value.endswith('%'):
            try:
                percent = float(value[:-1]) / 100.0
                return percent * reference_size
            except ValueError:
                return 0
        
        # Viewport units
        if value.endswith('vw'):
            try:
                vw = float(value[:-2]) / 100.0
                return vw * self.viewport_width
            except ValueError:
                return 0
        
        if value.endswith('vh'):
            try:
                vh = float(value[:-2]) / 100.0
                return vh * self.viewport_height
            except ValueError:
                return 0
        
        # Em units (approximate)
        if value.endswith('em'):
            try:
                em = float(value[:-2])
                return em * 16  # Assume 16px base font size
            except ValueError:
                return 0
        
        # Try to parse as direct number
        try:
            return float(value)
        except ValueError:
            return 0
    
    def _parse_box_dimensions(self, computed_style: Dict[str, str], property_prefix: str) -> Tuple[float, float, float, float]:
        """Parse box model dimensions (padding, margin, etc.)."""
        top = self._parse_dimension(computed_style.get(f'{property_prefix}-top', '0'))
        right = self._parse_dimension(computed_style.get(f'{property_prefix}-right', '0'))
        bottom = self._parse_dimension(computed_style.get(f'{property_prefix}-bottom', '0'))
        left = self._parse_dimension(computed_style.get(f'{property_prefix}-left', '0'))
        
        # Handle shorthand property
        shorthand = computed_style.get(property_prefix)
        if shorthand:
            values = shorthand.split()
            if len(values) == 1:
                # All sides same
                dim = self._parse_dimension(values[0])
                return (dim, dim, dim, dim)
            elif len(values) == 2:
                # Top/bottom, left/right
                tb = self._parse_dimension(values[0])
                lr = self._parse_dimension(values[1])
                return (tb, lr, tb, lr)
            elif len(values) == 4:
                # Top, right, bottom, left
                return tuple(self._parse_dimension(v) for v in values)
        
        return (top, right, bottom, left)


class EnhancedHtml2PptxConverter(Html2PptxConverter):
    """Enhanced HTML to PPTX converter with advanced layout analysis."""
    
    def __init__(self):
        super().__init__()
        self.css_parser = AdvancedCSSParser()
        self.element_bounds_cache = {}
        
        # Enhanced slide dimensions for better precision
        self.slide_width = Inches(13.33)  # 16:9 ratio
        self.slide_height = Inches(7.5)
        
        # Scaling factor to fit HTML viewport to slide
        self.scale_x = 13.33 / 19.2  # Slide inches / viewport width in inches (1920px @ 100dpi)
        self.scale_y = 7.5 / 10.8    # Slide inches / viewport height in inches (1080px @ 100dpi)
    
    def convert(self, input_path: str, output_path: str) -> None:
        """Enhanced conversion with advanced layout analysis."""
        
        print(f"Starting enhanced HTML to PPTX conversion...")
        
        try:
            self.base_path = os.path.dirname(os.path.abspath(input_path))
            
            # Read and parse HTML
            with open(input_path, 'r', encoding='utf-8') as f:
                html_content = f.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract and process CSS
            css_styles = self._extract_css_styles(soup)
            self._preprocess_css(css_styles)
            
            # Find slides with enhanced detection
            slides = self._detect_slides_enhanced(soup)
            
            if not slides:
                raise ConversionError("No slides found in HTML presentation")
            
            print(f"Detected {len(slides)} slides")
            
            # Create PowerPoint presentation
            prs = Presentation()
            
            # Remove default slide
            if len(prs.slides) > 0:
                slide_part = prs.slides._sldIdLst[0]
                prs.part.drop_rel(slide_part.rId)
                del prs.slides._sldIdLst[0]
            
            # Process each slide with enhanced analysis
            for slide_idx, slide_element in enumerate(slides):
                print(f"Processing slide {slide_idx + 1}/{len(slides)} with enhanced layout analysis...")
                
                try:
                    self._process_slide_enhanced(prs, slide_element, css_styles, slide_idx)
                except Exception as slide_error:
                    print(f"Warning: Enhanced processing failed for slide {slide_idx + 1}: {slide_error}")
                    print("Falling back to basic processing...")
                    
                    # Fallback to parent class method
                    try:
                        self._process_slide(prs, slide_element, css_styles, slide_idx)
                    except Exception as fallback_error:
                        print(f"Fallback also failed: {fallback_error}")
                        self._create_fallback_slide(prs, slide_idx, str(slide_error))
            
            # Set enhanced metadata
            self._set_enhanced_metadata(prs, soup)
            
            # Save with high quality settings
            prs.save(output_path)
            
            print(f"Successfully converted '{input_path}' to '{output_path}' with enhanced layout preservation")
            
        except Exception as e:
            raise ConversionError(
                f"Enhanced HTML to PPTX conversion failed: {e}",
                source_format="html",
                target_format="pptx",
                suggestions=[
                    "Install browser automation: pip install playwright",
                    "Ensure HTML has proper slide structure",
                    "Check CSS and layout complexity"
                ]
            )
    
    def _detect_slides_enhanced(self, soup: BeautifulSoup) -> List[Tag]:
        """Enhanced slide detection supporting multiple presentation formats."""
        
        # Try multiple slide detection patterns
        slide_selectors = [
            'div.slide',           # Standard slide class
            'section.slide',       # Section-based slides
            '.reveal .slides section',  # Reveal.js
            '.impress .step',      # Impress.js  
            '.deck-container .slide',   # Deck.js
            'article.slide',       # Article-based
            '[data-slide]',        # Data attribute based
        ]
        
        for selector in slide_selectors:
            try:
                slides = soup.select(selector)
                if slides:
                    print(f"Detected slides using selector: {selector}")
                    return slides
            except Exception:
                continue
        
        # Fallback to any element with 'slide' in class name
        potential_slides = soup.find_all(lambda tag: tag.get('class') and 
                                        any('slide' in cls.lower() for cls in tag.get('class', [])))
        
        if potential_slides:
            print(f"Using fallback slide detection: found {len(potential_slides)} potential slides")
            return potential_slides
        
        return []
    
    def _process_slide_enhanced(self, prs: Presentation, slide_element: Tag, css_styles: str, slide_idx: int):
        """Process slide with enhanced layout analysis and positioning."""
        
        # Add blank slide for precise control
        blank_slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Analyze slide layout
        layout_type = self.css_parser.detect_layout_type(slide_element)
        
        print(f"  Slide {slide_idx + 1}: Detected layout type: {layout_type.value}")
        
        # Calculate slide bounds
        slide_bounds = ElementBounds(0, 0, self.css_parser.viewport_width, self.css_parser.viewport_height)
        
        # Process elements based on layout type
        if layout_type == LayoutType.FLEX:
            self._process_flexbox_slide(slide, slide_element, css_styles, slide_bounds)
        elif layout_type == LayoutType.GRID:
            self._process_grid_slide(slide, slide_element, css_styles, slide_bounds)
        else:
            # Process with enhanced positioning analysis
            self._process_positioned_slide(slide, slide_element, css_styles, slide_bounds)
    
    def _process_flexbox_slide(self, slide, slide_element: Tag, css_styles: str, slide_bounds: ElementBounds):
        """Process slide with flexbox layout."""
        
        flex_properties = self.css_parser.analyze_flexbox_layout(slide_element)
        print(f"    Flexbox properties: {flex_properties}")
        
        # Get child elements
        children = [child for child in slide_element.children if isinstance(child, Tag)]
        
        # Calculate positions based on flexbox properties
        direction = flex_properties['direction']
        justify_content = flex_properties['justify_content']
        align_items = flex_properties['align_items']
        gap = flex_properties['gap']
        
        if direction in ['row', 'row-reverse']:
            # Horizontal layout
            available_width = slide_bounds.width
            child_width = (available_width - gap * (len(children) - 1)) / len(children) if children else 0
            
            for i, child in enumerate(children):
                x_pos = i * (child_width + gap)
                if direction == 'row-reverse':
                    x_pos = available_width - x_pos - child_width
                
                y_pos = self._calculate_flex_cross_position(slide_bounds.height, align_items)
                
                child_bounds = ElementBounds(x_pos, y_pos, child_width, slide_bounds.height * 0.8)
                self._create_positioned_element(slide, child, child_bounds, css_styles)
        
        else:
            # Vertical layout (column, column-reverse)
            available_height = slide_bounds.height
            child_height = (available_height - gap * (len(children) - 1)) / len(children) if children else 0
            
            for i, child in enumerate(children):
                y_pos = i * (child_height + gap)
                if direction == 'column-reverse':
                    y_pos = available_height - y_pos - child_height
                
                x_pos = self._calculate_flex_cross_position(slide_bounds.width, justify_content)
                
                child_bounds = ElementBounds(x_pos, y_pos, slide_bounds.width * 0.8, child_height)
                self._create_positioned_element(slide, child, child_bounds, css_styles)
    
    def _process_grid_slide(self, slide, slide_element: Tag, css_styles: str, slide_bounds: ElementBounds):
        """Process slide with CSS Grid layout."""
        
        grid_properties = self.css_parser.analyze_grid_layout(slide_element)
        print(f"    Grid properties: {grid_properties}")
        
        # Parse grid template
        columns = self._parse_grid_template(grid_properties['template_columns'], slide_bounds.width)
        rows = self._parse_grid_template(grid_properties['template_rows'], slide_bounds.height)
        
        # Get child elements
        children = [child for child in slide_element.children if isinstance(child, Tag)]
        
        # Position elements in grid
        for i, child in enumerate(children):
            # Calculate grid position (simplified - assumes auto-placement)
            if columns and rows:
                col = i % len(columns)
                row = i // len(columns)
                
                if row < len(rows):
                    x_pos = sum(columns[:col]) + col * grid_properties['column_gap']
                    y_pos = sum(rows[:row]) + row * grid_properties['row_gap'] 
                    width = columns[col]
                    height = rows[row]
                    
                    child_bounds = ElementBounds(x_pos, y_pos, width, height)
                    self._create_positioned_element(slide, child, child_bounds, css_styles)
    
    def _process_positioned_slide(self, slide, slide_element: Tag, css_styles: str, slide_bounds: ElementBounds):
        """Process slide with positioned elements (absolute, relative, static)."""
        
        # Find all significant child elements
        elements_to_process = []
        
        # Collect elements by priority (headers, content blocks, images, etc.)
        headers = slide_element.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
        content_blocks = slide_element.find_all(['div', 'section', 'article'], 
                                               class_=lambda x: x and any(cls in str(x).lower() 
                                               for cls in ['content', 'block', 'card', 'box', 'panel']))
        paragraphs = slide_element.find_all(['p'])
        lists = slide_element.find_all(['ul', 'ol'])
        tables = slide_element.find_all(['table'])
        images = slide_element.find_all(['img'])
        
        # Process by priority
        current_y = 50  # Start position
        
        # Headers first
        for header in headers:
            bounds = self._calculate_element_bounds(header, slide_bounds, current_y)
            self._create_positioned_element(slide, header, bounds, css_styles)
            current_y += bounds.height + 20
        
        # Content blocks
        for block in content_blocks:
            bounds = self._calculate_element_bounds(block, slide_bounds, current_y)
            self._create_positioned_element(slide, block, bounds, css_styles)
            current_y += bounds.height + 20
        
        # Other elements
        for element_list in [paragraphs, lists, tables, images]:
            for element in element_list:
                if current_y < slide_bounds.height - 100:  # Leave some margin
                    bounds = self._calculate_element_bounds(element, slide_bounds, current_y)
                    self._create_positioned_element(slide, element, bounds, css_styles)
                    current_y += bounds.height + 15
    
    def _create_positioned_element(self, slide, element: Tag, bounds: ElementBounds, css_styles: str):
        """Create PowerPoint element with precise positioning."""
        
        element_style = self.css_parser.parse_comprehensive_styles(element)
        
        # Convert bounds to PowerPoint coordinates
        left, top, width, height = bounds.to_inches(self.scale_x)
        
        # Ensure minimum sizes
        width = max(width, Inches(0.5))
        height = max(height, Inches(0.3))
        
        # Ensure element fits on slide
        if left + width > self.slide_width:
            width = self.slide_width - left
        if top + height > self.slide_height:
            height = self.slide_height - top
        
        # Create appropriate PowerPoint element
        if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
            self._create_positioned_heading(slide, element, left, top, width, height, element_style)
        
        elif element.name == 'p':
            self._create_positioned_paragraph(slide, element, left, top, width, height, element_style)
        
        elif element.name == 'img':
            self._create_positioned_image(slide, element, left, top, width, height, element_style)
        
        elif element.name in ['ul', 'ol']:
            self._create_positioned_list(slide, element, left, top, width, height, element_style)
        
        elif element.name == 'table':
            self._create_positioned_table(slide, element, left, top, width, height, element_style)
        
        else:
            # Generic text box for div and other elements
            self._create_positioned_text_box(slide, element, left, top, width, height, element_style)
    
    def _create_positioned_heading(self, slide, element: Tag, left, top, width, height, style: ElementStyle):
        """Create precisely positioned heading."""
        
        text_content = element.get_text(strip=True)
        if not text_content:
            return
        
        # Create text box
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Set text
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text_content
        
        # Apply styling
        font = paragraph.font
        
        # Font size based on heading level
        level = int(element.name[1]) if element.name[1:].isdigit() else 1
        if style.font_size:
            font.size = Pt(style.font_size)
        else:
            font.size = Pt(max(44 - (level - 1) * 6, 18))
        
        font.bold = (style.font_weight in ['bold', 'bolder', '700', '800', '900']) if style.font_weight else True
        
        if style.color:
            font.color.rgb = RGBColor(*style.color)
        
        # Text alignment
        if style.text_align == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif style.text_align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.LEFT
        
        # Apply background
        self._apply_shape_background(text_box, style)
    
    def _create_positioned_paragraph(self, slide, element: Tag, left, top, width, height, style: ElementStyle):
        """Create precisely positioned paragraph."""
        
        text_content = element.get_text(strip=True)
        if not text_content:
            return
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        paragraph = text_frame.paragraphs[0]
        paragraph.text = text_content
        
        # Apply styling
        font = paragraph.font
        
        if style.font_size:
            font.size = Pt(style.font_size)
        else:
            font.size = Pt(18)
        
        if style.font_weight in ['bold', 'bolder', '700', '800', '900']:
            font.bold = True
        
        if style.font_style == 'italic':
            font.italic = True
        
        if style.color:
            font.color.rgb = RGBColor(*style.color)
        
        # Text alignment
        if style.text_align == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif style.text_align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT
        elif style.text_align == 'justify':
            paragraph.alignment = PP_ALIGN.JUSTIFY
        
        self._apply_shape_background(text_box, style)
    
    def _create_positioned_text_box(self, slide, element: Tag, left, top, width, height, style: ElementStyle):
        """Create precisely positioned text box for complex content."""
        
        # Extract structured content
        text_content = self._extract_structured_text(element)
        if not text_content:
            return
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.clear()
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        
        # Handle multi-paragraph content
        lines = text_content.split('\\n')
        for i, line in enumerate(lines[:10]):  # Limit lines
            if i == 0:
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()
            
            paragraph.text = line.strip()
            
            # Apply styling
            font = paragraph.font
            font.size = Pt(style.font_size or 14)
            
            if style.color:
                font.color.rgb = RGBColor(*style.color)
        
        self._apply_shape_background(text_box, style)
    
    def _create_positioned_image(self, slide, element: Tag, left, top, width, height, style: ElementStyle):
        """Create precisely positioned image."""
        
        src = element.get('src')
        if not src:
            return
        
        try:
            if src.startswith('data:'):
                # Handle data URLs
                import base64
                from io import BytesIO
                
                header, data = src.split(',', 1)
                image_data = base64.b64decode(data)
                image_stream = BytesIO(image_data)
                slide.shapes.add_picture(image_stream, left, top, width, height)
                
            else:
                # Handle file paths
                image_path = src
                if not src.startswith(('http://', 'https://')):
                    image_path = os.path.join(self.base_path, src)
                
                if os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, left, top, width, height)
                    
        except Exception as e:
            print(f"Failed to add image {src}: {e}")
            # Create placeholder
            placeholder = slide.shapes.add_textbox(left, top, width, height)
            placeholder.text_frame.text = f"[Image: {src}]"
    
    def _apply_shape_background(self, shape, style: ElementStyle):
        """Apply background styling to shape."""
        
        if style.background_color:
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*style.background_color)
        
        elif style.background_gradient:
            # Apply gradient with proper PowerPoint API usage
            gradient = style.background_gradient
            if len(gradient['colors']) >= 2:
                try:
                    fill = shape.fill
                    
                    # Use the correct PowerPoint gradient API sequence
                    fill.gradient()
                    
                    # Set gradient angle
                    angle = gradient.get('angle', 0)
                    # Convert CSS angle (0=right, 90=down) to PowerPoint angle (0=right, 90=down)
                    fill.gradient_angle = angle
                    
                    # Get gradient stops and set colors
                    if hasattr(fill, 'gradient_stops') and hasattr(fill.gradient_stops, '__len__'):
                        stops = fill.gradient_stops
                        
                        # Clear existing stops by setting them to our colors
                        colors = gradient['colors']
                        for i in range(min(len(stops), len(colors))):
                            try:
                                stops[i].color.rgb = RGBColor(*colors[i])
                                # Set position (0.0 to 1.0)
                                stops[i].position = i / (len(colors) - 1) if len(colors) > 1 else 0.0
                            except Exception as stop_error:
                                print(f"Warning: Could not set gradient stop {i}: {stop_error}")
                                
                        # If we have more colors than stops, add additional stops
                        for i in range(len(stops), len(colors)):
                            try:
                                # Some PowerPoint APIs support adding new stops
                                if hasattr(stops, 'insert_gradient_stop'):
                                    position = i / (len(colors) - 1) if len(colors) > 1 else 0.0
                                    new_stop = stops.insert_gradient_stop(RGBColor(*colors[i]), position)
                            except Exception:
                                break  # Can't add more stops
                    else:
                        # Fallback: use simplified two-color gradient
                        try:
                            # Some versions use fore_color and back_color for gradients
                            fill.fore_color.rgb = RGBColor(*gradient['colors'][0])
                            if hasattr(fill, 'back_color'):
                                fill.back_color.rgb = RGBColor(*gradient['colors'][-1])
                        except Exception as color_error:
                            print(f"Warning: Could not set gradient colors: {color_error}")
                        
                except Exception as gradient_error:
                    print(f"Warning: Gradient application failed: {gradient_error}")
                    # Fallback to solid color using first gradient color
                    try:
                        fill = shape.fill
                        fill.solid()
                        fill.fore_color.rgb = RGBColor(*gradient['colors'][0])
                    except Exception as fallback_error:
                        print(f"Warning: Fallback color application also failed: {fallback_error}")
    
    def _calculate_element_bounds(self, element: Tag, slide_bounds: ElementBounds, suggested_y: float) -> ElementBounds:
        """Calculate element bounds with content-aware sizing."""
        
        # Parse explicit positioning from CSS
        bounds = self.css_parser.parse_element_bounds(element, slide_bounds)
        
        # If no explicit positioning, use suggested layout
        if bounds.x == 0 and bounds.y == 0:
            bounds.x = 50  # Default margin
            bounds.y = suggested_y
        
        # Content-aware width calculation
        if bounds.width <= 100:  # Default or small width
            text_length = len(element.get_text(strip=True))
            if text_length > 0:
                # Estimate width based on content
                estimated_width = min(text_length * 8, slide_bounds.width - 100)  # 8px per char estimate
                bounds.width = max(estimated_width, 200)
        
        # Content-aware height calculation  
        if bounds.height <= 50:  # Default or small height
            # Count content elements for height estimation
            lines = element.get_text().count('\\n') + 1
            if element.find(['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                bounds.height = max(lines * 25, 60)  # Headers need more space
            else:
                bounds.height = max(lines * 20, 40)  # Regular content
        
        return bounds
    
    def _extract_structured_text(self, element: Tag) -> str:
        """Extract structured text content preserving some formatting."""
        
        text_parts = []
        
        for child in element.descendants:
            if child.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                text_parts.append(f"\\n{child.get_text(strip=True)}\\n")
            elif child.name == 'p':
                text_parts.append(child.get_text(strip=True))
            elif child.name in ['li']:
                text_parts.append(f"• {child.get_text(strip=True)}")
            elif child.name == 'br':
                text_parts.append('\\n')
            elif hasattr(child, 'string') and child.string:
                text_parts.append(child.string.strip())
        
        return '\\n'.join(filter(None, text_parts))
    
    def _calculate_flex_cross_position(self, available_space: float, align_property: str) -> float:
        """Calculate cross-axis position for flexbox alignment."""
        
        if align_property == 'center':
            return available_space * 0.5
        elif align_property == 'flex-end':
            return available_space * 0.8
        else:  # flex-start, stretch, baseline
            return available_space * 0.1
    
    def _parse_grid_template(self, template: str, available_space: float) -> List[float]:
        """Parse CSS grid template into sizes."""
        
        if not template or template == 'none':
            return []
        
        # Handle common patterns
        if template == '1fr':
            return [available_space]
        elif 'fr' in template:
            # Parse fractional units
            parts = template.split()
            total_fr = sum(float(p.replace('fr', '')) for p in parts if 'fr' in p)
            if total_fr > 0:
                return [available_space * (float(p.replace('fr', '')) / total_fr) for p in parts if 'fr' in p]
        
        # Parse pixel and percentage values
        parts = template.split()
        sizes = []
        for part in parts:
            if part.endswith('px'):
                try:
                    sizes.append(float(part[:-2]))
                except ValueError:
                    pass
            elif part.endswith('%'):
                try:
                    percent = float(part[:-1]) / 100.0
                    sizes.append(available_space * percent)
                except ValueError:
                    pass
        
        return sizes if sizes else [available_space / 3] * 3  # Default 3 columns
    
    def _preprocess_css(self, css_styles: str):
        """Preprocess CSS for better parsing."""
        # Extract and cache important CSS rules
        # This is a simplified version - a full implementation would parse CSS properly
        
        # Extract viewport dimensions if specified
        viewport_match = re.search(r'width:\s*(\d+)px', css_styles)
        if viewport_match:
            self.css_parser.viewport_width = int(viewport_match.group(1))
        
        height_match = re.search(r'height:\s*(\d+)px', css_styles)
        if height_match:
            self.css_parser.viewport_height = int(height_match.group(1))
    
    def _set_enhanced_metadata(self, prs: Presentation, soup: BeautifulSoup):
        """Set enhanced presentation metadata."""
        
        # Call parent method first
        self._set_presentation_metadata(prs, soup)
        
        # Add additional metadata
        # Could extract more sophisticated metadata from HTML structure
        
        # Set slide size to 16:9 widescreen
        prs.slide_width = Emu(12192000)   # 13.33 inches
        prs.slide_height = Emu(6858000)   # 7.5 inches