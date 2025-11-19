"""
PPTX (PowerPoint) to various format converters.
"""

import sys
import os
from typing import Optional, List, Dict, Any

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
from core.ast import ASTNode, NodeType, StyleInfo, DocumentMetadata, create_document, create_paragraph, create_heading
from core.exceptions import ConversionError, DependencyError


class Pptx2MdConverter:
    """Converts PPTX files to Markdown format."""

    def parse_pptx2ast(self, input_path: str) -> ASTNode:
        """Parse PPTX file and convert to AST representation."""
        
        try:
            from pptx import Presentation
        except ImportError as e:
            raise DependencyError(
                "python-pptx is required for PPTX conversion",
                missing_dependency="python-pptx"
            ) from e

        # Load presentation
        prs = Presentation(input_path)
        doc = create_document()

        # Extract metadata
        if prs.core_properties:
            metadata = DocumentMetadata()
            if prs.core_properties.title:
                metadata.title = prs.core_properties.title
            if prs.core_properties.author:
                metadata.author = prs.core_properties.author
            if prs.core_properties.subject:
                metadata.subject = prs.core_properties.subject
            doc.metadata = metadata

        # Process slides
        for slide_idx, slide in enumerate(prs.slides):
            # Add slide header
            slide_header = create_heading(f"Slide {slide_idx + 1}", level=1)
            doc.add_child(slide_header)

            # Process shapes in slide
            for shape in slide.shapes:
                ast_node = self._process_pptx_shape(shape)
                if ast_node:
                    doc.add_child(ast_node)
            
            # Add slide separator
            if slide_idx < len(prs.slides) - 1:
                doc.add_child(ASTNode(NodeType.LINE_BREAK))

        return doc

    def _process_pptx_shape(self, shape) -> Optional[ASTNode]:
        """Process a PowerPoint shape and convert to AST node."""
        
        # Text boxes and title shapes
        if hasattr(shape, 'text_frame') and shape.text_frame:
            text_content = self._extract_text_from_textframe(shape.text_frame)
            if text_content.strip():
                # Check if it's a title (common PowerPoint pattern)
                if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                    ph_type = shape.placeholder_format.type
                    # Title placeholders (1 = title, 2 = subtitle)
                    if ph_type in [1, 2]:
                        level = 2 if ph_type == 1 else 3
                        return create_heading(text_content, level)
                
                return create_paragraph(text_content)

        # Tables
        elif hasattr(shape, 'table'):
            return self._process_pptx_table(shape.table)

        # Images
        elif shape.shape_type == 13:  # Picture type
            # Extract image info (limited without binary processing)
            image_node = ASTNode(NodeType.IMAGE)
            image_node.attributes = {
                'src': f'slide_image_{id(shape)}',
                'alt': 'Slide Image'
            }
            return image_node

        # Charts
        elif hasattr(shape, 'chart'):
            # Represent chart as paragraph with description
            return create_paragraph(f"[Chart: {shape.chart.chart_title.text_frame.text if shape.chart.chart_title else 'Untitled Chart'}]")

        return None

    def _extract_text_from_textframe(self, text_frame) -> str:
        """Extract text content from PowerPoint text frame."""
        
        text_parts = []
        
        for paragraph in text_frame.paragraphs:
            para_text = []
            
            for run in paragraph.runs:
                if run.text:
                    # Handle basic formatting
                    text = run.text
                    if run.font.bold:
                        text = f"**{text}**"
                    if run.font.italic:
                        text = f"*{text}*"
                    para_text.append(text)
            
            if para_text:
                text_parts.append(''.join(para_text))

        return '\\n'.join(text_parts)

    def _process_pptx_table(self, table) -> ASTNode:
        """Process PowerPoint table and convert to AST table node."""
        
        table_node = ASTNode(NodeType.TABLE)
        
        for row_idx, row in enumerate(table.rows):
            row_node = ASTNode(NodeType.TABLE_ROW)
            
            for cell in row.cells:
                cell_text = cell.text if cell.text else ''
                # First row is often headers in presentations
                cell_type = NodeType.TABLE_HEADER if row_idx == 0 else NodeType.TABLE_CELL
                cell_node = ASTNode(cell_type, content=cell_text)
                row_node.add_child(cell_node)
            
            table_node.add_child(row_node)

        return table_node

    def ast2md(self, ast_root: ASTNode, output_path: str) -> None:
        """Convert AST to Markdown format."""
        from core.mapper import ast2md
        ast2md(ast_root, output_path)

    def convert(self, input_path: str, output_path: str) -> None:
        """Convert PPTX to Markdown with enhanced Obsidian features."""
        
        try:
            # Use the enhanced Obsidian converter
            from .pptx2md_obsidian import Pptx2MdObsidianConverter
            obsidian_converter = Pptx2MdObsidianConverter()
            obsidian_converter.convert(input_path, output_path)
            
        except Exception as e:
            # Fallback to basic conversion
            try:
                ast_root = self.parse_pptx2ast(input_path)
                self.ast2md(ast_root, output_path)
                print(f"Successfully converted '{input_path}' to '{output_path}' (basic mode)")
            except Exception as fallback_error:
                raise ConversionError(
                    f"PPTX to Markdown conversion failed: Enhanced: {e}, Basic: {fallback_error}",
                    source_format="pptx",
                    target_format="md",
                    suggestions=[
                        "Install python-pptx: pip install python-pptx",
                        "Check that the PPTX file is valid and not corrupted",
                        "Ensure output directory has write permissions"
                    ]
                )


class Pptx2TxtConverter:
    """Converts PPTX files to plain text format."""

    def parse_pptx2ast(self, input_path: str) -> ASTNode:
        """Parse PPTX file and convert to AST representation."""
        pptx_converter = Pptx2MdConverter()
        return pptx_converter.parse_pptx2ast(input_path)

    def ast2txt(self, ast_root: ASTNode, output_path: str) -> None:
        """Convert AST to plain text format."""
        from core.mapper import ast2txt
        ast2txt(ast_root, output_path)

    def convert(self, input_path: str, output_path: str) -> None:
        """Convert PPTX to plain text."""
        
        try:
            # Method 1: AST-based conversion
            ast_root = self.parse_pptx2ast(input_path)
            self.ast2txt(ast_root, output_path)
            print(f"Successfully converted '{input_path}' to '{output_path}' (AST method)")
            
        except Exception as ast_error:
            try:
                # Method 2: Direct text extraction
                self._convert_direct(input_path, output_path)
                print(f"Successfully converted '{input_path}' to '{output_path}' (direct method)")
                
            except Exception as direct_error:
                raise ConversionError(
                    f"PPTX to TXT conversion failed. AST: {ast_error}, Direct: {direct_error}",
                    source_format="pptx",
                    target_format="txt"
                )

    def _convert_direct(self, input_path: str, output_path: str) -> None:
        """Direct PPTX to text extraction."""
        try:
            from pptx import Presentation
        except ImportError as e:
            raise DependencyError(
                "python-pptx is required for PPTX conversion",
                missing_dependency="python-pptx"
            ) from e

        prs = Presentation(input_path)
        text_content = []

        for slide_idx, slide in enumerate(prs.slides):
            text_content.append(f"=== SLIDE {slide_idx + 1} ===\\n")
            
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        para_text = ''.join(run.text for run in paragraph.runs)
                        if para_text.strip():
                            text_content.append(para_text)
                
                elif hasattr(shape, 'table'):
                    # Extract table as plain text
                    text_content.append("\\n[TABLE]")
                    for row in shape.table.rows:
                        row_text = ' | '.join(cell.text for cell in row.cells)
                        text_content.append(row_text)
                    text_content.append("[/TABLE]\\n")
            
            text_content.append('')  # Add blank line between slides

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\\n'.join(text_content))


class Pptx2PdfConverter:
    """Converts PPTX files to PDF format."""

    def parse_pptx2ast(self, input_path: str) -> ASTNode:
        """Parse PPTX file and convert to AST representation."""
        pptx_converter = Pptx2MdConverter()
        return pptx_converter.parse_pptx2ast(input_path)

    def convert(self, input_path: str, output_path: str) -> None:
        """Convert PPTX to PDF."""
        
        try:
            # Method 1: LibreOffice conversion (best quality)
            self._convert_with_libreoffice(input_path, output_path)
            print(f"Successfully converted '{input_path}' to '{output_path}' (LibreOffice)")
            
        except Exception as lo_error:
            try:
                # Method 2: Convert to HTML then PDF
                self._convert_via_html(input_path, output_path)
                print(f"Successfully converted '{input_path}' to '{output_path}' (HTML method)")
                
            except Exception as html_error:
                raise ConversionError(
                    f"PPTX to PDF conversion failed. LibreOffice: {lo_error}, HTML: {html_error}",
                    source_format="pptx",
                    target_format="pdf",
                    suggestions=[
                        "Install LibreOffice for best PPTX to PDF conversion",
                        "Install HTML to PDF dependencies as fallback"
                    ]
                )

    def _convert_with_libreoffice(self, input_path: str, output_path: str) -> None:
        """Convert using LibreOffice."""
        import subprocess
        import shutil
        
        if not shutil.which("libreoffice"):
            raise ConversionError("LibreOffice not found in PATH")
        
        output_dir = os.path.dirname(os.path.abspath(output_path))
        
        result = subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", output_dir, input_path
        ], capture_output=True, text=True)
        
        if result.returncode != 0:
            raise ConversionError(f"LibreOffice conversion failed: {result.stderr}")
        
        # Rename generated file if needed
        input_name = os.path.splitext(os.path.basename(input_path))[0]
        generated_pdf = os.path.join(output_dir, f"{input_name}.pdf")
        
        if generated_pdf != output_path and os.path.exists(generated_pdf):
            os.rename(generated_pdf, output_path)

    def _convert_via_html(self, input_path: str, output_path: str) -> None:
        """Convert PPTX to PDF via HTML intermediate."""
        import tempfile
        
        # Convert to HTML first
        with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as temp_html:
            ast_root = self.parse_pptx2ast(input_path)
            from core.mapper import ast2html
            ast2html(ast_root, temp_html.name)
            
            try:
                from .html2pdf import Html2PdfConverter
                pdf_converter = Html2PdfConverter()
                pdf_converter.convert(temp_html.name, output_path)
            finally:
                os.unlink(temp_html.name)


class Pptx2HtmlConverter:
    """Converts PPTX files to HTML format."""

    def parse_pptx2ast(self, input_path: str) -> ASTNode:
        """Parse PPTX file and convert to AST representation."""
        pptx_converter = Pptx2MdConverter()
        return pptx_converter.parse_pptx2ast(input_path)

    def ast2html(self, ast_root: ASTNode, output_path: str) -> None:
        """Convert AST to HTML format."""
        from core.mapper import ast2html
        ast2html(ast_root, output_path)

    def convert(self, input_path: str, output_path: str) -> None:
        """Convert PPTX to HTML."""
        
        try:
            # AST-based conversion
            ast_root = self.parse_pptx2ast(input_path)
            self.ast2html(ast_root, output_path)
            print(f"Successfully converted '{input_path}' to '{output_path}'")
            
        except Exception as e:
            raise ConversionError(
                f"PPTX to HTML conversion failed: {e}",
                source_format="pptx",
                target_format="html"
            )


# Reverse converters (other formats to PPTX)

class Md2PptxConverter:
    """Converts Markdown files to PPTX format."""

    def parse_md2ast(self, input_path: str) -> ASTNode:
        """Parse Markdown to AST representation."""
        from .md2html import Md2HtmlConverter
        md_converter = Md2HtmlConverter()
        return md_converter.parse_md2ast(input_path)

    def ast2pptx(self, ast_root: ASTNode, output_path: str) -> None:
        """Convert AST to PPTX format."""
        
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError as e:
            raise DependencyError(
                "python-pptx is required for PPTX creation",
                missing_dependency="python-pptx"
            ) from e

        prs = Presentation()
        
        # Set metadata if available
        if ast_root.metadata:
            if ast_root.metadata.title:
                prs.core_properties.title = ast_root.metadata.title
            if ast_root.metadata.author:
                prs.core_properties.author = ast_root.metadata.author

        current_slide = None
        
        for node in ast_root.children:
            if node.type == NodeType.HEADING:
                # Create new slide for each heading
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                current_slide = prs.slides.add_slide(slide_layout)
                
                title = current_slide.shapes.title
                title.text = node.content or ''
                
            elif node.type == NodeType.PARAGRAPH and current_slide:
                # Add content to current slide
                content = current_slide.placeholders[1]  # Content placeholder
                if hasattr(content, 'text_frame'):
                    p = content.text_frame.add_paragraph()
                    p.text = node.content or ''
                    
            elif node.type == NodeType.TABLE and current_slide:
                # Add table to slide (simplified)
                rows = len(node.find_children(NodeType.TABLE_ROW))
                if rows > 0:
                    first_row = node.find_children(NodeType.TABLE_ROW)[0]
                    cols = len([cell for cell in first_row.children 
                              if cell.type in [NodeType.TABLE_CELL, NodeType.TABLE_HEADER]])
                    
                    if rows > 0 and cols > 0:
                        # Add table shape
                        left = Inches(1.0)
                        top = Inches(3.0)
                        width = Inches(8.0)
                        height = Inches(2.0)
                        
                        table = current_slide.shapes.add_table(rows, cols, left, top, width, height).table
                        
                        # Fill table data
                        for row_idx, row_node in enumerate(node.find_children(NodeType.TABLE_ROW)):
                            for col_idx, cell_node in enumerate([cell for cell in row_node.children 
                                                               if cell.type in [NodeType.TABLE_CELL, NodeType.TABLE_HEADER]]):
                                if row_idx < len(table.rows) and col_idx < len(table.rows[row_idx].cells):
                                    table.rows[row_idx].cells[col_idx].text = cell_node.content or ''

        # If no slides were created, add a default title slide
        if len(prs.slides) == 0:
            slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = ast_root.metadata.title if ast_root.metadata and ast_root.metadata.title else "Converted Document"

        prs.save(output_path)

    def convert(self, input_path: str, output_path: str) -> None:
        """Convert Markdown to PPTX."""
        
        try:
            # Method 1: AST-based conversion
            ast_root = self.parse_md2ast(input_path)
            self.ast2pptx(ast_root, output_path)
            print(f"Successfully converted '{input_path}' to '{output_path}' (AST method)")
            
        except Exception as ast_error:
            try:
                # Method 2: Pandoc fallback
                self._convert_with_pandoc(input_path, output_path)
                print(f"Successfully converted '{input_path}' to '{output_path}' (pandoc fallback)")
                
            except Exception as pandoc_error:
                raise ConversionError(
                    f"Markdown to PPTX conversion failed. AST: {ast_error}, Pandoc: {pandoc_error}",
                    source_format="md",
                    target_format="pptx"
                )

    def _convert_with_pandoc(self, input_path: str, output_path: str) -> None:
        """Fallback conversion using pandoc."""
        try:
            import pypandoc
            pypandoc.convert_file(input_path, 'pptx', outputfile=output_path)
        except ImportError as e:
            raise DependencyError(
                "pypandoc is required for pandoc fallback conversion",
                missing_dependency="pypandoc"
            ) from e


class Txt2PptxConverter:
    """Converts plain text files to PPTX format."""

    def parse_txt2ast(self, input_path: str) -> ASTNode:
        """Parse plain text to AST representation."""
        from .txt2md import Txt2MdConverter
        txt_converter = Txt2MdConverter()
        return txt_converter.parse_txt2ast(input_path)

    def ast2pptx(self, ast_root: ASTNode, output_path: str) -> None:
        """Convert AST to PPTX format."""
        md_converter = Md2PptxConverter()
        md_converter.ast2pptx(ast_root, output_path)

    def convert(self, input_path: str, output_path: str) -> None:
        """Convert plain text to PPTX."""
        
        try:
            # AST-based conversion
            ast_root = self.parse_txt2ast(input_path)
            self.ast2pptx(ast_root, output_path)
            print(f"Successfully converted '{input_path}' to '{output_path}'")
            
        except Exception as e:
            raise ConversionError(
                f"TXT to PPTX conversion failed: {e}",
                source_format="txt",
                target_format="pptx"
            )