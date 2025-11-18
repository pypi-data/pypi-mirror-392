#!/usr/bin/env python3
"""MCP server for creating and managing PowerPoint presentations using python-pptx."""

import asyncio
import logging
from pathlib import Path
from typing import Any, Optional
import base64
import io

from mcp.server import Server
from mcp.types import (
    Tool,
    TextContent,
    ImageContent,
    EmbeddedResource,
)
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from PIL import Image as PILImage

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("powerpoint-mcp")

# Initialize MCP server
app = Server("powerpoint-mcp")


@app.list_tools()
async def list_tools() -> list[Tool]:
    """List available PowerPoint tools."""
    return [
        Tool(
            name="create_presentation",
            description="Create a new PowerPoint presentation with optional title slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path where the presentation will be saved (e.g., 'presentation.pptx')",
                    },
                    "title": {
                        "type": "string",
                        "description": "Title for the presentation (optional)",
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Subtitle for the title slide (optional)",
                    },
                },
                "required": ["filepath"],
            },
        ),
        Tool(
            name="add_title_slide",
            description="Add a title slide to an existing presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                    "title": {
                        "type": "string",
                        "description": "Title text",
                    },
                    "subtitle": {
                        "type": "string",
                        "description": "Subtitle text (optional)",
                    },
                },
                "required": ["filepath", "title"],
            },
        ),
        Tool(
            name="add_text_slide",
            description="Add a slide with title and content (bullet points or paragraphs)",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title",
                    },
                    "content": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of bullet points or content lines",
                    },
                },
                "required": ["filepath", "title", "content"],
            },
        ),
        Tool(
            name="add_image_slide",
            description="Add a slide with title and image",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                    "title": {
                        "type": "string",
                        "description": "Slide title",
                    },
                    "image_path": {
                        "type": "string",
                        "description": "Path to the image file to insert",
                    },
                    "left": {
                        "type": "number",
                        "description": "Left position in inches (optional, default: 1.5)",
                    },
                    "top": {
                        "type": "number",
                        "description": "Top position in inches (optional, default: 2.5)",
                    },
                    "width": {
                        "type": "number",
                        "description": "Image width in inches (optional, default: 6)",
                    },
                },
                "required": ["filepath", "title", "image_path"],
            },
        ),
        Tool(
            name="add_blank_slide",
            description="Add a blank slide to the presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                },
                "required": ["filepath"],
            },
        ),
        Tool(
            name="add_textbox",
            description="Add a textbox to the last slide with custom positioning and formatting",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                    "text": {
                        "type": "string",
                        "description": "Text content",
                    },
                    "left": {
                        "type": "number",
                        "description": "Left position in inches",
                    },
                    "top": {
                        "type": "number",
                        "description": "Top position in inches",
                    },
                    "width": {
                        "type": "number",
                        "description": "Width in inches",
                    },
                    "height": {
                        "type": "number",
                        "description": "Height in inches",
                    },
                    "font_size": {
                        "type": "number",
                        "description": "Font size in points (optional, default: 18)",
                    },
                    "bold": {
                        "type": "boolean",
                        "description": "Make text bold (optional, default: false)",
                    },
                },
                "required": ["filepath", "text", "left", "top", "width", "height"],
            },
        ),
        Tool(
            name="add_shape",
            description="Add a shape (rectangle, circle, etc.) to the last slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                    "shape_type": {
                        "type": "string",
                        "description": "Type of shape: 'rectangle', 'oval', 'rounded_rectangle', 'triangle'",
                        "enum": ["rectangle", "oval", "rounded_rectangle", "triangle"],
                    },
                    "left": {
                        "type": "number",
                        "description": "Left position in inches",
                    },
                    "top": {
                        "type": "number",
                        "description": "Top position in inches",
                    },
                    "width": {
                        "type": "number",
                        "description": "Width in inches",
                    },
                    "height": {
                        "type": "number",
                        "description": "Height in inches",
                    },
                    "fill_color": {
                        "type": "string",
                        "description": "Fill color in hex format (e.g., '#FF0000') (optional)",
                    },
                },
                "required": ["filepath", "shape_type", "left", "top", "width", "height"],
            },
        ),
        Tool(
            name="add_table",
            description="Add a table to the last slide",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                    "rows": {
                        "type": "number",
                        "description": "Number of rows",
                    },
                    "cols": {
                        "type": "number",
                        "description": "Number of columns",
                    },
                    "data": {
                        "type": "array",
                        "items": {
                            "type": "array",
                            "items": {"type": "string"},
                        },
                        "description": "2D array of table data (rows x cols)",
                    },
                    "left": {
                        "type": "number",
                        "description": "Left position in inches (optional, default: 1)",
                    },
                    "top": {
                        "type": "number",
                        "description": "Top position in inches (optional, default: 2)",
                    },
                    "width": {
                        "type": "number",
                        "description": "Table width in inches (optional, default: 8)",
                    },
                    "height": {
                        "type": "number",
                        "description": "Table height in inches (optional, default: 4)",
                    },
                },
                "required": ["filepath", "rows", "cols", "data"],
            },
        ),
        Tool(
            name="get_slide_count",
            description="Get the number of slides in a presentation",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                },
                "required": ["filepath"],
            },
        ),
        Tool(
            name="get_presentation_info",
            description="Get information about a presentation (title, slide count, etc.)",
            inputSchema={
                "type": "object",
                "properties": {
                    "filepath": {
                        "type": "string",
                        "description": "Path to the PowerPoint file",
                    },
                },
                "required": ["filepath"],
            },
        ),
    ]


def hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert hex color to RGB tuple."""
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


@app.call_tool()
async def call_tool(name: str, arguments: Any) -> list[TextContent]:
    """Handle tool calls for PowerPoint operations."""
    try:
        if name == "create_presentation":
            filepath = arguments["filepath"]
            title = arguments.get("title")
            subtitle = arguments.get("subtitle")

            prs = Presentation()

            # Add title slide if title is provided
            if title:
                slide_layout = prs.slide_layouts[0]  # Title slide layout
                slide = prs.slides.add_slide(slide_layout)

                title_shape = slide.shapes.title
                title_shape.text = title

                if subtitle and len(slide.placeholders) > 1:
                    subtitle_shape = slide.placeholders[1]
                    subtitle_shape.text = subtitle

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Successfully created presentation at {filepath}"
                    + (f" with title '{title}'" if title else ""),
                )
            ]

        elif name == "add_title_slide":
            filepath = arguments["filepath"]
            title = arguments["title"]
            subtitle = arguments.get("subtitle", "")

            prs = Presentation(filepath)
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = title

            if subtitle and len(slide.placeholders) > 1:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Added title slide '{title}' to {filepath}",
                )
            ]

        elif name == "add_text_slide":
            filepath = arguments["filepath"]
            title = arguments["title"]
            content = arguments["content"]

            prs = Presentation(filepath)
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = title

            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()

            for item in content:
                p = text_frame.add_paragraph()
                p.text = item
                p.level = 0

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Added text slide '{title}' with {len(content)} content items to {filepath}",
                )
            ]

        elif name == "add_image_slide":
            filepath = arguments["filepath"]
            title = arguments["title"]
            image_path = arguments["image_path"]
            left = arguments.get("left", 1.5)
            top = arguments.get("top", 2.5)
            width = arguments.get("width", 6)

            if not Path(image_path).exists():
                return [
                    TextContent(
                        type="text",
                        text=f"Error: Image file not found at {image_path}",
                    )
                ]

            prs = Presentation(filepath)
            slide_layout = prs.slide_layouts[5]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            # Add title
            title_shape = slide.shapes.title
            title_shape.text = title

            # Add image
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width),
            )

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Added image slide '{title}' with image from {image_path} to {filepath}",
                )
            ]

        elif name == "add_blank_slide":
            filepath = arguments["filepath"]

            prs = Presentation(filepath)
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Added blank slide to {filepath}",
                )
            ]

        elif name == "add_textbox":
            filepath = arguments["filepath"]
            text = arguments["text"]
            left = arguments["left"]
            top = arguments["top"]
            width = arguments["width"]
            height = arguments["height"]
            font_size = arguments.get("font_size", 18)
            bold = arguments.get("bold", False)

            prs = Presentation(filepath)

            if len(prs.slides) == 0:
                return [
                    TextContent(
                        type="text",
                        text=f"Error: No slides in presentation. Add a slide first.",
                    )
                ]

            slide = prs.slides[-1]  # Get last slide

            textbox = slide.shapes.add_textbox(
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            text_frame = textbox.text_frame
            text_frame.text = text

            # Format text
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)
                    run.font.bold = bold

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Added textbox to last slide in {filepath}",
                )
            ]

        elif name == "add_shape":
            filepath = arguments["filepath"]
            shape_type = arguments["shape_type"]
            left = arguments["left"]
            top = arguments["top"]
            width = arguments["width"]
            height = arguments["height"]
            fill_color = arguments.get("fill_color")

            prs = Presentation(filepath)

            if len(prs.slides) == 0:
                return [
                    TextContent(
                        type="text",
                        text=f"Error: No slides in presentation. Add a slide first.",
                    )
                ]

            slide = prs.slides[-1]

            # Map shape type to MSO_SHAPE constant
            shape_map = {
                "rectangle": MSO_SHAPE.RECTANGLE,
                "oval": MSO_SHAPE.OVAL,
                "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
                "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            }

            shape = slide.shapes.add_shape(
                shape_map[shape_type],
                Inches(left),
                Inches(top),
                Inches(width),
                Inches(height),
            )

            # Set fill color if provided
            if fill_color:
                rgb = hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(*rgb)

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Added {shape_type} shape to last slide in {filepath}",
                )
            ]

        elif name == "add_table":
            filepath = arguments["filepath"]
            rows = arguments["rows"]
            cols = arguments["cols"]
            data = arguments["data"]
            left = arguments.get("left", 1)
            top = arguments.get("top", 2)
            width = arguments.get("width", 8)
            height = arguments.get("height", 4)

            if len(data) != rows or any(len(row) != cols for row in data):
                return [
                    TextContent(
                        type="text",
                        text=f"Error: Data dimensions ({len(data)}x{len(data[0]) if data else 0}) don't match specified rows ({rows}) and cols ({cols})",
                    )
                ]

            prs = Presentation(filepath)

            if len(prs.slides) == 0:
                return [
                    TextContent(
                        type="text",
                        text=f"Error: No slides in presentation. Add a slide first.",
                    )
                ]

            slide = prs.slides[-1]

            table = slide.shapes.add_table(
                rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
            ).table

            # Fill table with data
            for i, row_data in enumerate(data):
                for j, cell_data in enumerate(row_data):
                    table.cell(i, j).text = str(cell_data)

            prs.save(filepath)

            return [
                TextContent(
                    type="text",
                    text=f"Added {rows}x{cols} table to last slide in {filepath}",
                )
            ]

        elif name == "get_slide_count":
            filepath = arguments["filepath"]

            if not Path(filepath).exists():
                return [
                    TextContent(
                        type="text",
                        text=f"Error: File not found at {filepath}",
                    )
                ]

            prs = Presentation(filepath)
            slide_count = len(prs.slides)

            return [
                TextContent(
                    type="text",
                    text=f"Presentation has {slide_count} slide(s)",
                )
            ]

        elif name == "get_presentation_info":
            filepath = arguments["filepath"]

            if not Path(filepath).exists():
                return [
                    TextContent(
                        type="text",
                        text=f"Error: File not found at {filepath}",
                    )
                ]

            prs = Presentation(filepath)
            slide_count = len(prs.slides)

            # Get core properties
            core_props = prs.core_properties

            info = f"""Presentation Information:
- File: {filepath}
- Slide count: {slide_count}
- Title: {core_props.title or 'Not set'}
- Author: {core_props.author or 'Not set'}
- Subject: {core_props.subject or 'Not set'}
- Created: {core_props.created or 'Unknown'}
- Modified: {core_props.modified or 'Unknown'}
"""

            return [TextContent(type="text", text=info)]

        else:
            return [
                TextContent(
                    type="text",
                    text=f"Unknown tool: {name}",
                )
            ]

    except Exception as e:
        logger.error(f"Error executing {name}: {e}", exc_info=True)
        return [
            TextContent(
                type="text",
                text=f"Error executing {name}: {str(e)}",
            )
        ]


async def main():
    """Run the MCP server."""
    from mcp.server.stdio import stdio_server

    async with stdio_server() as (read_stream, write_stream):
        logger.info("PowerPoint MCP server starting...")
        await app.run(
            read_stream,
            write_stream,
            app.create_initialization_options(),
        )


if __name__ == "__main__":
    asyncio.run(main())
